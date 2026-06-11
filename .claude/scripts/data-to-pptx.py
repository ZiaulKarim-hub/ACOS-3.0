#!/usr/bin/env python3
"""
ACOS PPTX Generation Engine — Component Library + Layout Engine

Generates institutional-quality PPTX presentations from structured data.
Bypasses HTML intermediate — builds slides natively via python-pptx.

Usage:
    python3 data-to-pptx.py <data_yaml> <design_spec_yaml> [--template template.pptx] [-o output.pptx]

Input:
    data_yaml        — verified-data.yaml (loan data with provenance)
    design_spec_yaml — design-spec.yaml (colors, fonts, layout from Phase 1)
    --template       — optional blank .pptx with theme/masters
    -o               — output path (default: output.pptx)

All text uses shape.text_frame (never floating textboxes).
All text_frames have vertical_anchor set explicitly.
All margins set to 0 or intentional values.
Font selected by content role, not hardcoded.
"""
import argparse
import re
import sys
import yaml
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# =============================================================================
# CONSTANTS (defaults — overridden by design spec when provided)
# =============================================================================

DEFAULT_SLIDE_WIDTH = Inches(13.333)   # 16:9 widescreen
DEFAULT_SLIDE_HEIGHT = Inches(7.5)

# Default margins (from slide edge)
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.4)

# Content area (after header/footer)
HEADER_HEIGHT = Inches(0.8)
FOOTER_HEIGHT = Inches(0.45)

# Font roles — auto-selected by content type
FONT_ROLES = {
    "number":  "IBM Plex Mono",  # financial figures, percentages, dates
    "display": "IBM Plex Sans",  # titles, headings, prose
    "label":   "IBM Plex Sans",  # labels, captions, metadata
}

# Default colors (overridden by design spec)
DEFAULT_COLORS = {
    "primary":    RGBColor(0x46, 0x5D, 0x53),  # sage-80 (OKOA primary)
    "secondary":  RGBColor(0x10, 0x44, 0x73),  # navy (OKOA secondary)
    "accent":     RGBColor(0xFF, 0x79, 0x5E),  # coral (OKOA accent)
    "bg_dark":    RGBColor(0x46, 0x5D, 0x53),  # sage-80
    "bg_light":   RGBColor(0xF4, 0xF4, 0xF4),  # ink-05
    "text_dark":  RGBColor(0x16, 0x16, 0x16),  # ink
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),  # white
    "positive":   RGBColor(0x24, 0xA1, 0x48),  # OKOA green
    "caution":    RGBColor(0xF1, 0xC2, 0x1B),  # OKOA amber
    "negative":   RGBColor(0xDA, 0x1E, 0x28),  # OKOA red
    "border":     RGBColor(0xE0, 0xE0, 0xE0),  # ink-10
}


# =============================================================================
# DESIGN SPEC LOADER
# =============================================================================

def load_design_spec(path):
    """Load design spec YAML and merge with defaults."""
    spec = {}
    if path and Path(path).exists():
        with open(path) as f:
            spec = yaml.safe_load(f) or {}
        if not isinstance(spec, dict):
            print(f"WARNING: design spec is not a YAML mapping — using defaults", file=sys.stderr)
            spec = {}

    # --- Slide dimensions (override from spec or use defaults) ---
    spec["_slide_width"] = DEFAULT_SLIDE_WIDTH
    spec["_slide_height"] = DEFAULT_SLIDE_HEIGHT
    if spec.get("slide_width"):
        try:
            spec["_slide_width"] = Inches(float(spec["slide_width"]))
        except (ValueError, TypeError):
            pass
    if spec.get("slide_height"):
        try:
            spec["_slide_height"] = Inches(float(spec["slide_height"]))
        except (ValueError, TypeError):
            pass

    # --- Colors (handle flat, nested, and missing-# formats) ---
    colors = dict(DEFAULT_COLORS)
    raw_colors = spec.get("colors", {})
    if isinstance(raw_colors, dict):
        # Handle nested structure: {palette: {primary: "#hex"}}
        flat_colors = {}
        for key, val in raw_colors.items():
            if isinstance(val, dict):
                # Nested: e.g. colors.palette.primary
                for sub_key, sub_val in val.items():
                    flat_colors[sub_key] = sub_val
            else:
                flat_colors[key] = val

        for key, hex_val in flat_colors.items():
            if isinstance(hex_val, str):
                hex_val = hex_val.lstrip("#")
                if len(hex_val) == 6:
                    try:
                        colors[key] = RGBColor(
                            int(hex_val[0:2], 16),
                            int(hex_val[2:4], 16),
                            int(hex_val[4:6], 16),
                        )
                    except (ValueError, IndexError):
                        pass  # keep default color
    spec["_colors"] = colors

    fonts = dict(FONT_ROLES)
    if "fonts" in spec:
        for role, name in spec["fonts"].items():
            if role in fonts:
                fonts[role] = name
    spec["_fonts"] = fonts

    return spec


def load_data(path):
    """Load verified data YAML with error handling."""
    if not path or not Path(path).exists():
        print(f"ERROR: data file not found: {path}", file=sys.stderr)
        sys.exit(2)
    with open(path) as f:
        try:
            data = yaml.safe_load(f)
        except yaml.YAMLError as e:
            print(f"ERROR: failed to parse YAML in {path}: {e}", file=sys.stderr)
            sys.exit(2)
    if not data:
        print(f"ERROR: data file is empty or invalid: {path}", file=sys.stderr)
        sys.exit(2)
    return data


# =============================================================================
# LAYOUT ENGINE
# =============================================================================

class ContentArea:
    """Defines usable slide area after header/footer."""
    def __init__(self, slide_width=None, slide_height=None,
                 header_h=HEADER_HEIGHT, footer_h=FOOTER_HEIGHT):
        sw = slide_width if slide_width is not None else DEFAULT_SLIDE_WIDTH
        sh = slide_height if slide_height is not None else DEFAULT_SLIDE_HEIGHT
        self.left = MARGIN_LEFT
        self.top = MARGIN_TOP + header_h
        self.width = sw - MARGIN_LEFT - MARGIN_RIGHT
        self.height = sh - MARGIN_TOP - MARGIN_BOTTOM - header_h - footer_h


class SlideGrid:
    """Computes cell positions from cols/rows/gap specs within a content area."""
    def __init__(self, area, cols=1, rows=1, h_gap=Inches(0.2), v_gap=Inches(0.2)):
        self.area = area
        self.cols = cols
        self.rows = rows
        self.h_gap = h_gap
        self.v_gap = v_gap
        self.cell_width = (area.width - h_gap * (cols - 1)) // cols
        self.cell_height = (area.height - v_gap * (rows - 1)) // rows
        self.cell_width = max(self.cell_width, Emu(1))
        self.cell_height = max(self.cell_height, Emu(1))

    def cell(self, col, row):
        """Return (left, top, width, height) for grid cell at (col, row), 0-indexed."""
        left = self.area.left + col * (self.cell_width + self.h_gap)
        top = self.area.top + row * (self.cell_height + self.v_gap)
        return left, top, self.cell_width, self.cell_height

    def span(self, col_start, row_start, col_span=1, row_span=1):
        """Return bounds for a multi-cell span."""
        left = self.area.left + col_start * (self.cell_width + self.h_gap)
        top = self.area.top + row_start * (self.cell_height + self.v_gap)
        width = col_span * self.cell_width + (col_span - 1) * self.h_gap
        height = row_span * self.cell_height + (row_span - 1) * self.v_gap
        return left, top, width, height


class CarbonGrid(SlideGrid):
    """16-column grid aligned to Carbon 2x Grid for consistent slide layouts.

    Provides column-span positioning that mirrors the CSS 16-column grid
    used in HTML/PDF output, ensuring visual consistency across formats.

    Usage:
        grid = CarbonGrid(area)
        left, top, width, height = grid.columns(start=0, span=8)  # left half
        left, top, width, height = grid.columns(start=8, span=8)  # right half
    """
    GUTTER = Inches(0.15)   # ~32px equivalent at slide scale
    COLUMNS = 16

    def __init__(self, area):
        super().__init__(area, cols=self.COLUMNS, rows=1,
                         h_gap=self.GUTTER, v_gap=Inches(0))

    def columns(self, start=0, span=16):
        """Return (left, top, width, height) for a column span.

        Args:
            start: 0-indexed column start (0-15)
            span: number of columns to span (1-16)

        Returns:
            (left, top, width, height) in EMU
        """
        return self.span(col_start=start, row_start=0,
                         col_span=span, row_span=1)


# =============================================================================
# TEXT HELPERS
# =============================================================================

def detect_font_role(text):
    """Auto-detect content type for font role selection.

    Rules (checked in priority order):
    1. "number" -- text IS primarily a numeric value (currency, %, ratio, date-only)
    2. "label" -- short uppercase headers, category labels, column headers
    3. "display" -- everything else (prose, titles, entity names, descriptions)

    Key principle: text that CONTAINS a number is NOT necessarily "number" role.
    "Village 1", "Scenario A -- 6-Month Hold", "& 1722 Mohawk LLC" are display, not number.
    """
    if text is None or text == "":
        return "label"
    s = str(text).strip()

    # 1. Pure numeric values: "$6,000,000", "14.00%", "50.00%", "$2,333", "9.75% . 19.5% Ann."
    #    Strip currency/percent/comma/spaces/punctuation -- if >=60% of remainder is digits, it's number
    cleaned = re.sub(r'[\$,%\s\u00b7\-\u2013+~\u00d7]', '', s)
    if cleaned and len(cleaned) <= 20:
        digit_ratio = sum(c.isdigit() or c == '.' for c in cleaned) / len(cleaned)
        if digit_ratio >= 0.6:
            return "number"

    # 2. Date-only values: "Oct 1, 2026", "Apr 4, 2024", "MARCH 2026"
    #    But NOT sentences like "Effective Apr 4, 2025" or "+180 days from Apr 3"
    if re.match(r'^(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+\d', s, re.IGNORECASE):
        if len(s.split()) <= 3:
            return "number"

    # 3. Short number+unit patterns: "6-12 Mo.", "~113 Acres", "2 / 5", "1st D/T"
    if re.match(r'^[~$]?\d[\d,.\-\u2013/\s]*\s*(Mo\.?|Acres?|Days?|p\.a\.|Ann\.?|D/T)?\s*$', s, re.IGNORECASE):
        return "number"

    # 4. Financial column/unit headers: "% OF LOAN", "PER DIEM"
    if re.match(r'^[%$#]\s', s) or s.upper() in ("% OF LOAN", "PER DIEM", "P.A."):
        return "number"

    # 5. UPPERCASE short labels: section headers, category labels
    if s.isupper() and len(s) < 50:
        return "label"

    # 6. Short non-numeric text: labels, captions
    if len(s) < 30 and '.' not in s and not any(c.isdigit() for c in s):
        return "label"

    # 7. Everything else: display (prose, titles, entity names, descriptions)
    return "display"


def set_text(tf, text, font_size=Pt(12), bold=False, color=None,
             alignment=PP_ALIGN.LEFT, font_name=None, spec=None):
    """Set text on an existing text_frame's first paragraph."""
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    elif spec:
        run.font.name = spec["_fonts"].get(detect_font_role(text), "Calibri")
    return p


def add_paragraph(tf, text, font_size=Pt(12), bold=False, color=None,
                  alignment=PP_ALIGN.LEFT, font_name=None, spec=None,
                  space_before=Pt(0), space_after=Pt(0)):
    """Add a new paragraph to an existing text_frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    elif spec:
        run.font.name = spec["_fonts"].get(detect_font_role(text), "Calibri")
    return p


def zero_margins(tf):
    """Set all internal margins to zero."""
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def set_margins(tf, left=Emu(0), right=Emu(0), top=Emu(0), bottom=Emu(0)):
    """Set explicit margins on a text_frame."""
    tf.margin_left = left
    tf.margin_right = right
    tf.margin_top = top
    tf.margin_bottom = bottom


def compute_height(font_size_pt, line_count, line_spacing=1.4):
    """Compute textbox height from content."""
    return Pt(font_size_pt * line_count * line_spacing)


def check_text_fit(shape):
    """Check if text fits within shape bounds after font changes. Returns True if OK."""
    if not shape.has_text_frame:
        return True
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            if not run.text.strip() or not run.font.size:
                continue
            font_size_pt = run.font.size / 12700  # EMU → pt: 1 pt = 12700 EMU
            is_mono = run.font.name in ("Courier New", "Consolas")
            char_width_pt = font_size_pt * (0.62 if is_mono else 0.48)
            text_width_pt = len(run.text) * char_width_pt
            box_width_pt = shape.width / 12700
            margin_pt = ((tf.margin_left or 0) + (tf.margin_right or 0)) / 12700
            available = box_width_pt - margin_pt
            if text_width_pt > available * 1.05:  # 5% tolerance
                return False
    return True


def auto_fit_font(shape, min_size_pt=7):
    """Reduce font size in 0.5pt increments until text fits, with a minimum."""
    if check_text_fit(shape):
        return False  # no change needed
    tf = shape.text_frame
    changed = False
    for p in tf.paragraphs:
        for run in p.runs:
            if not run.font.size or not run.text.strip():
                continue
            current_pt = run.font.size / 12700
            while current_pt > min_size_pt:
                if check_text_fit(shape):
                    break
                current_pt -= 0.5
                run.font.size = Pt(current_pt)
                changed = True
    if not check_text_fit(shape):
        print(f"WARNING: text cannot fit shape even at {min_size_pt}pt", file=sys.stderr)
    return changed


# =============================================================================
# SLIDE BOUNDARY CLAMPING
# =============================================================================

def clamp_to_slide(shape, slide_width, slide_height, footer_height=FOOTER_HEIGHT):
    """Ensure a shape does not extend beyond slide boundaries.

    Adjusts position/size if the shape overflows the slide area.
    Footer height is reserved at the bottom.
    """
    if shape.width <= 0 or shape.height <= 0:
        return

    max_bottom = slide_height - footer_height

    # Clamp right edge
    if shape.left + shape.width > slide_width:
        overflow = (shape.left + shape.width) - slide_width
        if overflow <= shape.width // 2:
            shape.width = slide_width - shape.left
        else:
            shape.left = slide_width - shape.width
            if shape.left < 0:
                shape.left = Emu(0)
                shape.width = slide_width

    # Clamp bottom edge
    if shape.top + shape.height > max_bottom:
        overflow = (shape.top + shape.height) - max_bottom
        if overflow <= shape.height // 2:
            shape.height = max_bottom - shape.top
        else:
            shape.top = max_bottom - shape.height
            if shape.top < 0:
                shape.top = Emu(0)
                shape.height = max_bottom

    # Clamp left edge (should not go negative)
    if shape.left < 0:
        shape.left = Emu(0)

    # Clamp top edge
    if shape.top < 0:
        shape.top = Emu(0)


# =============================================================================
# SHAPE HELPERS
# =============================================================================

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape with optional fill and border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    """Add a textbox with explicit anchor and zero margins."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    zero_margins(tf)
    try:
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    except Exception:
        pass
    # Set vertical anchor
    tf.vertical_anchor = anchor
    return txbox, tf


# =============================================================================
# COMPONENTS
# =============================================================================

def header_bar(slide, title, subtitle=None, spec=None):
    """Render a slide header bar across the top.

    Title and subtitle are placed in separate textboxes:
    title left-aligned at left margin, subtitle right-aligned at right side.
    Both vertically centered within the header bar.
    """
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    slide_width = spec.get("_slide_width", DEFAULT_SLIDE_WIDTH) if spec else DEFAULT_SLIDE_WIDTH

    # Background bar
    add_rect(slide, Emu(0), Emu(0), slide_width, HEADER_HEIGHT + MARGIN_TOP,
             fill_color=colors["bg_dark"])

    # Title textbox (left-aligned, takes left ~60% of width)
    title_width = Inches(7.0) if not subtitle else slide_width * 6 // 10
    txbox_title, tf_title = add_textbox(slide, MARGIN_LEFT, MARGIN_TOP,
                                        title_width, HEADER_HEIGHT,
                                        anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf_title, title, font_size=Pt(22), bold=True,
             color=colors["text_light"],
             font_name=spec["_fonts"]["display"] if spec else "Georgia",
             spec=spec)

    # Subtitle textbox (right-aligned, separate from title)
    if subtitle:
        sub_width = slide_width * 35 // 100
        sub_left = slide_width - MARGIN_RIGHT - sub_width
        txbox_sub, tf_sub = add_textbox(slide, sub_left, MARGIN_TOP,
                                        sub_width, HEADER_HEIGHT,
                                        anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf_sub, subtitle, font_size=Pt(12), bold=False,
                 color=colors["text_light"], alignment=PP_ALIGN.RIGHT,
                 font_name=spec["_fonts"]["label"] if spec else "Calibri",
                 spec=spec)


def footer_bar(slide, left_text="", right_text="", spec=None):
    """Render a slide footer bar across the bottom."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    slide_width = spec.get("_slide_width", DEFAULT_SLIDE_WIDTH) if spec else DEFAULT_SLIDE_WIDTH
    slide_height = spec.get("_slide_height", DEFAULT_SLIDE_HEIGHT) if spec else DEFAULT_SLIDE_HEIGHT

    footer_top = slide_height - FOOTER_HEIGHT - MARGIN_BOTTOM
    # Separator line
    add_rect(slide, MARGIN_LEFT, footer_top, slide_width - MARGIN_LEFT * 2,
             Pt(1), fill_color=colors["border"])
    # Left text
    if left_text:
        txbox, tf = add_textbox(slide, MARGIN_LEFT, footer_top + Pt(4),
                                Inches(5), FOOTER_HEIGHT - Pt(4),
                                anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, left_text, font_size=Pt(8), color=colors["text_dark"],
                 font_name=spec["_fonts"]["label"] if spec else "Calibri", spec=spec)
    # Right text
    if right_text:
        txbox, tf = add_textbox(slide, slide_width - MARGIN_RIGHT - Inches(3),
                                footer_top + Pt(4), Inches(3), FOOTER_HEIGHT - Pt(4),
                                anchor=MSO_ANCHOR.MIDDLE)
        set_text(tf, right_text, font_size=Pt(8), color=colors["text_dark"],
                 alignment=PP_ALIGN.RIGHT,
                 font_name=spec["_fonts"]["label"] if spec else "Calibri", spec=spec)


def section_header(slide, title, area, spec=None):
    """Render a section divider heading within the content area."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    # Accent bar
    add_rect(slide, area.left, area.top, Inches(0.08), Inches(0.35),
             fill_color=colors["accent"])
    # Title
    txbox, tf = add_textbox(slide, area.left + Inches(0.2), area.top,
                            area.width - Inches(0.2), Inches(0.35),
                            anchor=MSO_ANCHOR.MIDDLE)
    set_text(tf, title.upper(), font_size=Pt(11), bold=True,
             color=colors["primary"],
             font_name=spec["_fonts"]["label"] if spec else "Calibri", spec=spec)
    return Inches(0.45)  # height consumed


def metric_card(slide, left, top, width, height, label, value, sub_text=None,
                bg_color=None, spec=None):
    """Render a metric card using shape.text_frame with paragraphs (no floating textboxes)."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    if bg_color is None:
        bg_color = colors["bg_light"]

    shape = add_rect(slide, left, top, width, height, fill_color=bg_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    set_margins(tf, left=Inches(0.15), right=Inches(0.15),
                top=Inches(0.12), bottom=Inches(0.08))

    # Label paragraph
    p_label = tf.paragraphs[0]
    p_label.alignment = PP_ALIGN.LEFT
    p_label.space_after = Pt(2)
    run = p_label.runs[0] if p_label.runs else p_label.add_run()
    run.text = str(label).upper()
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = colors["secondary"]
    run.font.name = spec["_fonts"]["label"] if spec else "Calibri"

    # Value paragraph
    p_val = tf.add_paragraph()
    p_val.alignment = PP_ALIGN.LEFT
    p_val.space_before = Pt(2)
    p_val.space_after = Pt(2)
    run_v = p_val.add_run()
    run_v.text = str(value)
    run_v.font.size = Pt(18)
    run_v.font.bold = True
    run_v.font.color.rgb = colors["text_dark"]
    run_v.font.name = spec["_fonts"]["number"] if spec else "Courier New"

    # Sub-text paragraph (optional)
    if sub_text:
        p_sub = tf.add_paragraph()
        p_sub.alignment = PP_ALIGN.LEFT
        p_sub.space_before = Pt(1)
        run_s = p_sub.add_run()
        run_s.text = str(sub_text)
        run_s.font.size = Pt(8)
        run_s.font.color.rgb = colors["secondary"]
        run_s.font.name = spec["_fonts"]["label"] if spec else "Calibri"

    # Vertical anchor: middle if no sub-text, top otherwise
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if not sub_text else MSO_ANCHOR.TOP


def badge_pill(slide, left, top, width, height, text, bg_color=None, text_color=None, spec=None):
    """Render a badge/pill (e.g., recommendation status) with centered text."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    if bg_color is None:
        bg_color = colors["positive"]
    if text_color is None:
        text_color = colors["text_light"]

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    zero_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text).upper()
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = spec["_fonts"]["label"] if spec else "Calibri"


def data_table(slide, left, top, width, headers, rows, spec=None, col_widths=None):
    """Render a native table with styled header and alternating rows."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    row_height = Inches(0.35)
    table_height = row_height * num_rows
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = table_shape.table

    # Set column widths (only if count matches to avoid IndexError)
    if col_widths and len(col_widths) == num_cols:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        col_w = width // num_cols
        for i in range(num_cols):
            table.columns[i].width = col_w

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["bg_dark"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.text = str(header_text)
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = colors["text_light"]
        run.font.name = spec["_fonts"]["label"] if spec else "Calibri"
        cell.text_frame.margin_left = Inches(0.08)
        cell.text_frame.margin_top = Pt(2)
        cell.text_frame.margin_bottom = Pt(2)

    # Data rows with alternating backgrounds
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            # Alternating row colors
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors["bg_light"]
            else:
                cell.fill.background()

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            safe_text = str(cell_text) if cell_text is not None else ""
            run.text = safe_text
            # Use number font for numeric content
            role = detect_font_role(safe_text)
            run.font.name = spec["_fonts"].get(role, "Calibri") if spec else FONT_ROLES.get(role, "Calibri")
            run.font.size = Pt(9)
            run.font.color.rgb = colors["text_dark"]
            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_top = Pt(2)
            cell.text_frame.margin_bottom = Pt(2)

    return table_shape


def timeline_strip(slide, left, top, width, height, events, spec=None):
    """Render a horizontal timeline of date/event/detail cards."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    n = len(events)
    if n == 0:
        return

    card_gap = Inches(0.15)
    card_width = (width - card_gap * (n - 1)) // n
    card_height = height

    for i, event in enumerate(events):
        card_left = left + i * (card_width + card_gap)
        # Accent bar at top
        add_rect(slide, card_left, top, card_width, Pt(3),
                 fill_color=colors["accent"])
        # Card background
        shape = add_rect(slide, card_left, top + Pt(3), card_width,
                         card_height - Pt(3), fill_color=colors["bg_light"])
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        set_margins(tf, left=Inches(0.1), right=Inches(0.1),
                    top=Inches(0.08), bottom=Inches(0.05))
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # Date
        p_date = tf.paragraphs[0]
        p_date.alignment = PP_ALIGN.LEFT
        p_date.space_after = Pt(3)
        run = p_date.runs[0] if p_date.runs else p_date.add_run()
        run.text = str(event.get("date", ""))
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = colors["primary"]
        run.font.name = spec["_fonts"]["number"] if spec else "Courier New"

        # Event title
        p_title = tf.add_paragraph()
        p_title.alignment = PP_ALIGN.LEFT
        p_title.space_after = Pt(2)
        run_t = p_title.add_run()
        run_t.text = str(event.get("title", ""))
        run_t.font.size = Pt(9)
        run_t.font.bold = True
        run_t.font.color.rgb = colors["text_dark"]
        run_t.font.name = spec["_fonts"]["label"] if spec else "Calibri"

        # Detail (optional)
        if event.get("detail"):
            p_detail = tf.add_paragraph()
            p_detail.alignment = PP_ALIGN.LEFT
            run_d = p_detail.add_run()
            run_d.text = str(event["detail"])
            run_d.font.size = Pt(8)
            run_d.font.color.rgb = colors["secondary"]
            run_d.font.name = spec["_fonts"]["label"] if spec else "Calibri"


def scenario_box(slide, left, top, width, height, title, line_items,
                 total_label=None, total_value=None, accent_color=None, spec=None):
    """Render a scenario box: accent bar + line items + separator + totals."""
    colors = spec["_colors"] if spec else DEFAULT_COLORS
    if accent_color is None:
        accent_color = colors["accent"]

    # Accent bar on left
    add_rect(slide, left, top, Inches(0.06), height, fill_color=accent_color)

    # Main box
    shape = add_rect(slide, left + Inches(0.06), top, width - Inches(0.06), height,
                     fill_color=colors["bg_light"], border_color=colors["border"],
                     border_width=Pt(0.5))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    set_margins(tf, left=Inches(0.15), right=Inches(0.15),
                top=Inches(0.1), bottom=Inches(0.08))
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Title
    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    p_title.space_after = Pt(6)
    run = p_title.runs[0] if p_title.runs else p_title.add_run()
    run.text = str(title).upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = colors["primary"]
    run.font.name = spec["_fonts"]["label"] if spec else "Calibri"

    # Line items
    for item in line_items:
        label = item.get("label", "")
        value = item.get("value", "")
        line_text = f"{label}    {value}" if value else str(label)
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run_l = p.add_run()
        run_l.text = line_text
        run_l.font.size = Pt(9)
        run_l.font.color.rgb = colors["text_dark"]
        run_l.font.name = spec["_fonts"]["number"] if spec else "Courier New"

    # Separator + total
    if total_label and total_value:
        p_sep = tf.add_paragraph()
        p_sep.alignment = PP_ALIGN.LEFT
        p_sep.space_before = Pt(4)
        p_sep.space_after = Pt(2)
        run_sep = p_sep.add_run()
        run_sep.text = "\u2500" * 40
        run_sep.font.size = Pt(6)
        run_sep.font.color.rgb = colors["border"]

        total_text = f"{total_label}    {total_value}"
        p_total = tf.add_paragraph()
        p_total.alignment = PP_ALIGN.LEFT
        run_tot = p_total.add_run()
        run_tot.text = total_text
        run_tot.font.size = Pt(10)
        run_tot.font.bold = True
        run_tot.font.color.rgb = colors["primary"]
        run_tot.font.name = spec["_fonts"]["number"] if spec else "Courier New"


# =============================================================================
# COMPONENT DISPATCH (for build_content_slide)
# =============================================================================

COMPONENT_DISPATCH = {
    "metric_card": lambda slide, comp, area, spec: metric_card(
        slide,
        left=area.left + Inches(comp.get("x_offset", 0)),
        top=area.top + Inches(comp.get("y_offset", 0)),
        width=Inches(comp.get("width", 2.5)),
        height=Inches(comp.get("height", 1.2)),
        label=comp.get("label", ""),
        value=comp.get("value", ""),
        sub_text=comp.get("sub_text"),
        bg_color=None,
        spec=spec,
    ),
    "data_table": lambda slide, comp, area, spec: data_table(
        slide,
        left=area.left + Inches(comp.get("x_offset", 0)),
        top=area.top + Inches(comp.get("y_offset", 0)),
        width=Inches(comp.get("width", 12)),
        headers=comp.get("headers", []),
        rows=comp.get("rows", []),
        spec=spec,
    ),
    "timeline_strip": lambda slide, comp, area, spec: timeline_strip(
        slide,
        left=area.left + Inches(comp.get("x_offset", 0)),
        top=area.top + Inches(comp.get("y_offset", 0)),
        width=Inches(comp.get("width", 12)),
        height=Inches(comp.get("height", 1.5)),
        events=comp.get("events", []),
        spec=spec,
    ),
    "scenario_box": lambda slide, comp, area, spec: scenario_box(
        slide,
        left=area.left + Inches(comp.get("x_offset", 0)),
        top=area.top + Inches(comp.get("y_offset", 0)),
        width=Inches(comp.get("width", 5)),
        height=Inches(comp.get("height", 3)),
        title=comp.get("title", ""),
        line_items=comp.get("line_items", []),
        total_label=comp.get("total_label"),
        total_value=comp.get("total_value"),
        spec=spec,
    ),
    "badge_pill": lambda slide, comp, area, spec: badge_pill(
        slide,
        left=area.left + Inches(comp.get("x_offset", 0)),
        top=area.top + Inches(comp.get("y_offset", 0)),
        width=Inches(comp.get("width", 2.5)),
        height=Inches(comp.get("height", 0.4)),
        text=comp.get("text", ""),
        spec=spec,
    ),
    "section_header": lambda slide, comp, area, spec: section_header(
        slide,
        title=comp.get("title", ""),
        area=area,
        spec=spec,
    ),
}


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_cover_slide(prs, data, spec):
    """Build a title/cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    colors = spec["_colors"]
    slide_width = spec.get("_slide_width", DEFAULT_SLIDE_WIDTH)
    slide_height = spec.get("_slide_height", DEFAULT_SLIDE_HEIGHT)

    # Full-slide dark background
    add_rect(slide, Emu(0), Emu(0), slide_width, slide_height,
             fill_color=colors["bg_dark"])

    # Title
    txbox, tf = add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                            anchor=MSO_ANCHOR.BOTTOM)
    set_text(tf, data.get("title", "Presentation"), font_size=Pt(36), bold=True,
             color=colors["text_light"], font_name=spec["_fonts"]["display"], spec=spec)

    # Subtitle
    if data.get("subtitle"):
        txbox2, tf2 = add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                                  anchor=MSO_ANCHOR.TOP)
        set_text(tf2, data["subtitle"], font_size=Pt(16), bold=False,
                 color=colors["secondary"], font_name=spec["_fonts"]["label"], spec=spec)

    # Date / metadata
    if data.get("date"):
        txbox3, tf3 = add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
                                  anchor=MSO_ANCHOR.TOP)
        set_text(tf3, data["date"], font_size=Pt(12), bold=False,
                 color=colors["text_light"], font_name=spec["_fonts"]["number"], spec=spec)

    # Accent line
    add_rect(slide, Inches(1), Inches(3.3), Inches(2), Pt(3),
             fill_color=colors["accent"])

    return slide


def build_content_slide(prs, slide_spec, data, spec):
    """Build a generic content slide from a slide specification.

    Reads slide_spec["components"] and renders each one via COMPONENT_DISPATCH.
    If no components are specified, the slide is left with just header/footer
    (the Phase 3 agent will add content programmatically).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide_width = spec.get("_slide_width", DEFAULT_SLIDE_WIDTH)
    slide_height = spec.get("_slide_height", DEFAULT_SLIDE_HEIGHT)
    area = ContentArea(slide_width=slide_width, slide_height=slide_height)

    title = slide_spec.get("title", "")
    header_bar(slide, title, subtitle=slide_spec.get("subtitle"), spec=spec)
    footer_bar(slide, left_text=data.get("footer_left", ""),
               right_text=data.get("footer_right", ""), spec=spec)

    # Render components if specified
    components = slide_spec.get("components", [])
    for comp in components:
        comp_type = comp.get("type", "")
        handler = COMPONENT_DISPATCH.get(comp_type)
        if handler:
            result = handler(slide, comp, area, spec)
            # Clamp shapes placed by components to slide boundaries
            # result may be a shape, table_shape, or None
            if result is not None and hasattr(result, 'left'):
                clamp_to_slide(result, slide_width, slide_height, FOOTER_HEIGHT)

    return slide, area


# =============================================================================
# LOAN-DATA FORMAT ADAPTER
# =============================================================================

def auto_generate_slides_from_loan_data(data, spec):
    """Generate slide specs from loan-data.yaml format.

    When the PPTX pipeline receives a loan-data.yaml (with financial_figures,
    entities, data_by_section) instead of a pre-structured slides list, this
    function synthesizes a minimal slide deck that the Phase 3 agent can
    populate programmatically.
    """
    slides = []

    # Cover slide
    title = data.get("deal_name", data.get("property_name", "Loan Participation Offering"))
    slides.append({
        "type": "cover",
        "title": title,
        "subtitle": data.get("document_title", ""),
        "date": data.get("date", ""),
    })

    # Investment Overview slide from financial_figures
    if data.get("financial_figures") or data.get("entities"):
        slides.append({
            "type": "content",
            "title": "Investment Overview",
            "components": []  # Phase 3 agent fills these programmatically
        })

    # One slide per section from data_by_section
    if data.get("data_by_section"):
        for section_name, section_data in data["data_by_section"].items():
            slides.append({
                "type": "content",
                "title": section_name,
                "components": []
            })

    return slides


# =============================================================================
# MAIN ORCHESTRATOR
# =============================================================================

def generate_presentation(data_path, spec_path, template_path=None, output_path="output.pptx"):
    """Generate a complete PPTX from data and design spec."""
    data = load_data(data_path)
    spec = load_design_spec(spec_path)

    slide_width = spec.get("_slide_width", DEFAULT_SLIDE_WIDTH)
    slide_height = spec.get("_slide_height", DEFAULT_SLIDE_HEIGHT)

    # Create presentation from template or blank
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Always set slide dimensions from spec (overrides template defaults too)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Build slides from data structure
    slides_data = data.get("slides", [])
    if not slides_data:
        # Auto-generate from loan-data.yaml format (financial_figures, data_by_section, entities)
        if data.get("financial_figures") or data.get("data_by_section") or data.get("entities"):
            slides_data = auto_generate_slides_from_loan_data(data, spec)
        else:
            # Last resort: single cover slide.
            # Normalize loan-data keys (deal_name/property_name/document_title)
            # into the title/subtitle/date keys build_cover_slide expects —
            # mirrors the cover mapping in auto_generate_slides_from_loan_data.
            cover_spec = {
                "type": "cover",
                "title": data.get("deal_name",
                                  data.get("property_name", "Loan Participation Offering")),
                "subtitle": data.get("document_title", ""),
                "date": data.get("date", ""),
            }
            build_cover_slide(prs, cover_spec, spec)

    # Iterate over slides (works for both explicit and auto-generated slide lists)
    for slide_spec in slides_data:
        slide_type = slide_spec.get("type", "content")
        if slide_type == "cover":
            build_cover_slide(prs, slide_spec, spec)
        elif slide_type == "content":
            build_content_slide(prs, slide_spec, data, spec)
        else:
            print(f"WARNING: unrecognized slide type '{slide_type}' — skipped", file=sys.stderr)

    # Save — ensure output directory exists
    output_dir = Path(output_path).parent
    if not output_dir.exists():
        output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Generated: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="ACOS PPTX Generation Engine")
    parser.add_argument("data_yaml", help="Path to verified-data.yaml")
    parser.add_argument("design_spec_yaml", help="Path to design-spec.yaml")
    parser.add_argument("--template", help="Optional template.pptx with theme/masters")
    parser.add_argument("-o", "--output", default="output.pptx", help="Output path")
    args = parser.parse_args()

    generate_presentation(args.data_yaml, args.design_spec_yaml,
                          template_path=args.template, output_path=args.output)


if __name__ == "__main__":
    main()
