#!/usr/bin/env python3
"""
ACOS PPTX Validation Script

Automated post-generation checks for PPTX presentations.
Produces a YAML validation report compatible with Phase 4 format.

Usage:
    python3 validate-pptx.py <pptx_path> [<verified_data_path>] [<design_spec_path>]

Checks:
    1. Text overflow — every shape's text fits within bounds
    2. Boundary — no shape extends beyond slide edges
    3. Data integrity — every number matches verified-data.yaml
    4. Font compliance — correct font per content role
    5. Color compliance — all RGB values in design spec palette
    6. Anchor audit — all text_frames have vertical_anchor set
    7. Margin audit — no default margins (all explicit)
"""
import argparse
import re
import sys
import yaml
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.oxml.ns import qn
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


FONT_ROLES = {
    "number":  "IBM Plex Mono",  # financial figures, percentages, dates
    "display": "IBM Plex Sans",  # titles, headings, prose
    "label":   "IBM Plex Sans",  # labels, captions, metadata
}

# Default python-pptx margins (91440 EMU = 0.1 inch)
DEFAULT_MARGIN = 91440


def detect_font_role(text):
    """Auto-detect content type for font role selection.
    Must match the logic in data-to-pptx.py exactly."""
    if text is None or text == "":
        return "label"
    s = str(text).strip()

    # 1. Pure numeric values (>=60% digits after stripping currency/punctuation)
    cleaned = re.sub(r'[\$,%\s·\-–+~×]', '', s)
    if cleaned and len(cleaned) <= 20:
        digit_ratio = sum(c.isdigit() or c == '.' for c in cleaned) / len(cleaned)
        if digit_ratio >= 0.6:
            return "number"

    # 2. Date-only values (short: "Oct 1, 2026" but not "Effective Apr 4, 2025")
    if re.match(r'^(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+\d', s, re.IGNORECASE):
        if len(s.split()) <= 3:
            return "number"

    # 3. Short number+unit: "6–12 Mo.", "~113 Acres", "2 / 5", "1st D/T"
    if re.match(r'^[~$]?\d[\d,.\-–/\s]*\s*(Mo\.?|Acres?|Days?|p\.a\.|Ann\.?|D/T)?\s*$', s, re.IGNORECASE):
        return "number"

    # 4. Financial column/unit headers
    if re.match(r'^[%$#]\s', s) or s.upper() in ("% OF LOAN", "PER DIEM", "P.A."):
        return "number"

    # 5. UPPERCASE short labels
    if s.isupper() and len(s) < 50:
        return "label"

    # 6. Short non-numeric text
    if len(s) < 30 and '.' not in s and not any(c.isdigit() for c in s):
        return "label"

    # 7. Everything else: display
    return "display"


def load_yaml(path):
    if not path or not Path(path).exists():
        return {}
    try:
        with open(path) as f:
            data = yaml.safe_load(f)
            return data or {}
    except yaml.YAMLError as e:
        print(f"WARNING: could not parse YAML at {path}: {e}", file=sys.stderr)
        return {}


def extract_numbers(text):
    """Extract all numeric values from text for data integrity checks."""
    if not text:
        return []
    # Match numbers with optional currency, commas, decimals, percentages
    pattern = r'[\$]?\d[\d,]*\.?\d*[%x]?'
    return re.findall(pattern, str(text))


def normalize_number(s):
    """Normalize a number string for comparison."""
    return s.replace("$", "").replace(",", "").replace("%", "").replace("x", "").strip()


class PptxValidator:
    def __init__(self, pptx_path, data_path=None, spec_path=None):
        self.pptx_path = pptx_path
        self._broken = False
        self._broken_error = None
        try:
            self.prs = Presentation(pptx_path)
        except Exception as e:
            self._broken = True
            self._broken_error = str(e)
            self.prs = None
            self.findings = []
            return
        self._data_path_provided = bool(data_path)
        self.data = load_yaml(data_path) if data_path else {}
        self.spec = load_yaml(spec_path) if spec_path else {}
        self.findings = []
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        # Build color palette from spec (flatten nested dicts)
        self.palette = set()
        if "colors" in self.spec:
            color_data = self.spec["colors"]
            flat_colors = {}
            for key, val in color_data.items():
                if isinstance(val, dict):
                    for k2, v2 in val.items():
                        flat_colors[k2] = v2
                else:
                    flat_colors[key] = val
            for key, hex_val in flat_colors.items():
                if isinstance(hex_val, str):
                    hex_val = hex_val.lstrip("#")
                    if len(hex_val) == 6:
                        self.palette.add(hex_val.upper())

        # Build expected fonts from spec
        self.expected_fonts = dict(FONT_ROLES)
        if "fonts" in self.spec:
            for role, name in self.spec["fonts"].items():
                if role in self.expected_fonts:
                    self.expected_fonts[role] = name

        # Collect all numbers from verified data for integrity check
        self.data_numbers = set()
        self._collect_numbers(self.data)

    def _collect_numbers(self, obj, prefix=""):
        if isinstance(obj, dict):
            for k, v in obj.items():
                self._collect_numbers(v, f"{prefix}.{k}")
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                self._collect_numbers(v, f"{prefix}[{i}]")
        elif isinstance(obj, (int, float)):
            self.data_numbers.add(str(obj))
        elif isinstance(obj, str):
            for n in extract_numbers(obj):
                self.data_numbers.add(normalize_number(n))

    def add_finding(self, check, severity, slide_num, shape_name, message):
        self.findings.append({
            "check": check,
            "severity": severity,
            "slide": slide_num,
            "shape": shape_name,
            "message": message,
        })

    def check_text_overflow(self):
        """Check that text fits within shape bounds (height and width)."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                text = tf.text.strip()
                if not text:
                    continue

                # Placeholder shapes that inherit position/size from the slide
                # layout return None for these attrs; skip them rather than
                # crash on the arithmetic below (shape.height * 1.2, etc.).
                if shape.width is None or shape.height is None:
                    continue

                # Height check: estimate from font sizes and line count
                lines = text.split("\n")
                max_font_size = 12
                for p in tf.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            pts = run.font.size.pt
                            if pts > max_font_size:
                                max_font_size = pts

                estimated_height = Emu(int(max_font_size * len(lines) * 1.4 * 12700))
                if estimated_height > shape.height * 1.2:
                    self.add_finding(
                        "text_overflow", "warning", slide_idx,
                        shape.name,
                        f"Text may overflow vertically: {len(lines)} lines at "
                        f"~{max_font_size}pt in {shape.height} EMU height"
                    )

                # Width check: account for monospace being wider
                margin_emu = (tf.margin_left or 0) + (tf.margin_right or 0)
                available_width = shape.width - margin_emu
                if available_width <= 0:
                    continue
                for p in tf.paragraphs:
                    for run in p.runs:
                        if not run.text.strip() or not run.font.size:
                            continue
                        font_pt = run.font.size.pt
                        is_mono = run.font.name in ("Courier New", "Consolas")
                        char_width_emu = int(font_pt * (0.62 if is_mono else 0.48) * 12700)
                        text_width = len(run.text) * char_width_emu
                        if text_width > available_width * 1.05:
                            self.add_finding(
                                "text_overflow", "warning", slide_idx,
                                shape.name,
                                f"Text may overflow horizontally: '{run.text[:30]}...' "
                                f"({run.font.name}, {font_pt}pt) est. {text_width} EMU "
                                f"> available {available_width} EMU"
                            )

    def check_boundaries(self):
        """Check no shape extends beyond slide edges."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                # Placeholder shapes that inherit position/size from the slide
                # layout return None for left/top/width/height; computing
                # shape.left + shape.width on them raises TypeError and aborts
                # the whole validation run. Skip them (matches check-pptx-layout.py).
                if (shape.left is None or shape.top is None
                        or shape.width is None or shape.height is None):
                    continue
                right = shape.left + shape.width
                bottom = shape.top + shape.height
                if shape.left < 0:
                    self.add_finding("boundary", "error", slide_idx, shape.name,
                                     f"Shape extends beyond left edge: left={shape.left}")
                if shape.top < 0:
                    self.add_finding("boundary", "error", slide_idx, shape.name,
                                     f"Shape extends beyond top edge: top={shape.top}")
                if right > self.slide_width:
                    self.add_finding("boundary", "error", slide_idx, shape.name,
                                     f"Shape extends beyond right edge: right={right} > {self.slide_width}")
                if bottom > self.slide_height:
                    self.add_finding("boundary", "error", slide_idx, shape.name,
                                     f"Shape extends beyond bottom edge: bottom={bottom} > {self.slide_height}")

    def check_data_integrity(self):
        """Check numbers in slides match verified data."""
        if not self.data_numbers:
            if self._data_path_provided:
                self.add_finding("data_integrity", "warning", 0, "",
                    "Data integrity check skipped — verified data file was empty or unreadable")
            return

        slide_numbers = set()
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for num in extract_numbers(shape.text_frame.text):
                    normalized = normalize_number(num)
                    if normalized and len(normalized) > 1:  # skip single digits
                        slide_numbers.add(normalized)

        # Check for numbers in slides that don't appear in data
        unmatched = []
        for num in slide_numbers:
            if num not in self.data_numbers:
                # Check with some tolerance for rounding
                found = False
                try:
                    n = float(num)
                    for dn in self.data_numbers:
                        try:
                            if abs(float(dn) - n) < max(0.01 * abs(n), 0.005):
                                found = True
                                break
                        except ValueError:
                            continue
                except ValueError:
                    pass
                if not found:
                    unmatched.append(num)

        if unmatched:
            self.add_finding(
                "data_integrity", "warning", 0, "",
                f"{len(unmatched)} number(s) in presentation not found in verified data: "
                f"{', '.join(unmatched[:5])}{'...' if len(unmatched) > 5 else ''}"
            )

    def check_font_compliance(self):
        """Check fonts match expected roles."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if not run.font.name or not run.text.strip():
                            continue
                        expected_role = detect_font_role(run.text)
                        expected_font = self.expected_fonts.get(expected_role)
                        if expected_font and run.font.name != expected_font:
                            self.add_finding(
                                "font_compliance", "info", slide_idx, shape.name,
                                f"Font '{run.font.name}' used for {expected_role} content "
                                f"(expected '{expected_font}'): '{run.text[:30]}...'"
                            )

    def check_color_compliance(self):
        """Check colors are within the design spec palette."""
        if not self.palette:
            return

        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.color and run.font.color.rgb:
                            hex_color = str(run.font.color.rgb).upper()
                            if hex_color not in self.palette and hex_color not in ("000000", "FFFFFF"):
                                self.add_finding(
                                    "color_compliance", "info", slide_idx, shape.name,
                                    f"Color #{hex_color} not in design spec palette"
                                )

    def check_anchor_audit(self):
        """Check all text_frames have vertical_anchor set."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # Skip placeholder shapes from slide masters
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    continue
                if not shape.text_frame.text.strip():
                    continue
                bps = shape._element.findall('.//' + qn('a:bodyPr'))
                has_anchor = any(bp.get('anchor') is not None for bp in bps)
                if not has_anchor:
                    self.add_finding(
                        "anchor_audit", "warning", slide_idx, shape.name,
                        "text_frame has no vertical_anchor set"
                    )

    def check_margin_audit(self):
        """Check no text_frame has default (unset) margins."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # Skip placeholder shapes from slide masters (margins are
                # inherited from the layout, not spurious "default" findings).
                if getattr(shape, 'is_placeholder', False):
                    continue
                if not shape.text_frame.text.strip():
                    continue
                tf = shape.text_frame
                defaults = []
                if tf.margin_left == DEFAULT_MARGIN:
                    defaults.append("margin_left")
                if tf.margin_right == DEFAULT_MARGIN:
                    defaults.append("margin_right")
                if tf.margin_top == DEFAULT_MARGIN:
                    defaults.append("margin_top")
                if tf.margin_bottom == DEFAULT_MARGIN:
                    defaults.append("margin_bottom")
                if defaults:
                    self.add_finding(
                        "margin_audit", "info", slide_idx, shape.name,
                        f"Default margins detected: {', '.join(defaults)}"
                    )

    def validate(self):
        """Run all checks and return report."""
        if self._broken:
            return {
                "validation": {
                    "pptx_path": self.pptx_path,
                    "slide_count": 0,
                    "verdict": "FAIL",
                    "summary": {
                        "errors": 1,
                        "warnings": 0,
                        "info": 0,
                        "total": 1,
                    },
                    "findings": [{
                        "check": "file_integrity",
                        "severity": "error",
                        "slide": 0,
                        "shape": "",
                        "message": f"Corrupt or unreadable PPTX: {self._broken_error}",
                    }],
                }
            }
        self.check_boundaries()
        self.check_text_overflow()
        self.check_anchor_audit()
        self.check_margin_audit()
        self.check_font_compliance()
        self.check_color_compliance()
        self.check_data_integrity()

        errors = [f for f in self.findings if f["severity"] == "error"]
        warnings = [f for f in self.findings if f["severity"] == "warning"]
        infos = [f for f in self.findings if f["severity"] == "info"]

        report = {
            "validation": {
                "pptx_path": self.pptx_path,
                "slide_count": len(self.prs.slides),
                "verdict": "PASS" if not errors else "FAIL",
                "summary": {
                    "errors": len(errors),
                    "warnings": len(warnings),
                    "info": len(infos),
                    "total": len(self.findings),
                },
                "findings": self.findings,
            }
        }
        return report


def main():
    parser = argparse.ArgumentParser(description="ACOS PPTX Validation")
    parser.add_argument("pptx_path", help="Path to generated .pptx")
    parser.add_argument("verified_data_path", nargs="?", help="Path to verified-data.yaml")
    parser.add_argument("design_spec_path", nargs="?", help="Path to design-spec.yaml")
    parser.add_argument("-o", "--output", help="Write YAML report to file")
    args = parser.parse_args()

    validator = PptxValidator(args.pptx_path, args.verified_data_path, args.design_spec_path)
    report = validator.validate()

    report_yaml = yaml.dump(report, default_flow_style=False, sort_keys=False)

    if args.output:
        with open(args.output, "w") as f:
            f.write(report_yaml)
        print(f"Report written to: {args.output}")
    else:
        print(report_yaml)

    verdict = report["validation"]["verdict"]
    summary = report["validation"]["summary"]
    print(f"\nVerdict: {verdict} ({summary['errors']} errors, "
          f"{summary['warnings']} warnings, {summary['info']} info)", file=sys.stderr)
    sys.exit(0 if verdict == "PASS" else 1)


if __name__ == "__main__":
    main()
