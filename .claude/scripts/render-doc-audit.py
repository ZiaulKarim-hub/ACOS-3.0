#!/usr/bin/env python3
"""
ACOS Document Screenshot Renderer -- Universal Visual QA Tool

Renders any document (PPTX, PDF, DOCX) to PNG screenshots for visual QA.
For PPTX files, performs shape-level analysis including overflow detection,
overlap detection, and bounding box visualization.

Usage:
    python3 render-doc-audit.py <input-file> [--output-dir <dir>] [--dpi <200>] [--format pptx|pdf|docx|auto]

Dependencies:
    Required: python-pptx, Pillow (for PPTX rendering)
    Optional: pdf2image / poppler (for PDF), textutil or libreoffice (for DOCX)
"""

import argparse
import os
import shutil
import subprocess
import sys
import tempfile
import textwrap
from pathlib import Path

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

# ---------------------------------------------------------------------------
# Conditional imports -- graceful degradation
# ---------------------------------------------------------------------------

HAS_PPTX = False
HAS_PIL = False
HAS_PDF2IMAGE = False

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.oxml.ns import qn
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    pass

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    pass

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    pass


# ---------------------------------------------------------------------------
# Character width estimation table (EMU per character, calibrated)
# ---------------------------------------------------------------------------

CHAR_WIDTHS = {
    'Georgia': {
        7: 52000, 8: 58000, 9: 65000, 10: 72000, 11: 78000,
        12: 85000, 14: 100000, 16: 115000, 18: 130000,
        20: 145000, 24: 170000, 28: 200000, 32: 230000,
    },
    'Calibri': {
        7: 48000, 8: 54000, 9: 60000, 10: 67000, 11: 73000,
        12: 80000, 14: 93000, 16: 107000, 18: 120000, 20: 134000,
    },
    'Courier New': {
        16: 130000, 18: 146000, 20: 162000,
    },
}

# Fallback: 65000 EMU per character at Pt(9), scaled linearly
FALLBACK_CHAR_WIDTH_AT_9PT = 65000


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def estimate_char_width_emu(font_name, font_size_pt):
    """Return estimated character width in EMU for a given font and size."""
    font_name = font_name or 'Calibri'
    font_size_pt = font_size_pt or 9

    table = CHAR_WIDTHS.get(font_name)
    if table:
        # Exact match
        if font_size_pt in table:
            return table[font_size_pt]
        # Interpolate between nearest sizes
        sizes = sorted(table.keys())
        if font_size_pt < sizes[0]:
            return int(table[sizes[0]] * font_size_pt / sizes[0])
        if font_size_pt > sizes[-1]:
            return int(table[sizes[-1]] * font_size_pt / sizes[-1])
        # Find bracketing sizes
        for i in range(len(sizes) - 1):
            if sizes[i] <= font_size_pt <= sizes[i + 1]:
                lo, hi = sizes[i], sizes[i + 1]
                ratio = (font_size_pt - lo) / (hi - lo)
                return int(table[lo] + ratio * (table[hi] - table[lo]))

    # Fallback: scale from 9pt baseline
    return int(FALLBACK_CHAR_WIDTH_AT_9PT * font_size_pt / 9.0)


def emu_to_px(emu, dpi):
    """Convert EMU to pixels at a given DPI."""
    # 1 inch = 914400 EMU
    return int(emu * dpi / 914400)


def px_to_emu(px, dpi):
    """Convert pixels to EMU at a given DPI."""
    return int(px * 914400 / dpi)


def shapes_overlap(s1, s2):
    """Check if two shape bounding boxes overlap (given as (left, top, width, height) in EMU)."""
    l1, t1, w1, h1 = s1
    l2, t2, w2, h2 = s2

    # No overlap if one is entirely to the left/right/above/below the other
    if l1 + w1 <= l2 or l2 + w2 <= l1:
        return False
    if t1 + h1 <= t2 or t2 + h2 <= t1:
        return False
    return True


def extract_fill_color(shape):
    """Extract fill color from a shape. Returns (R, G, B) tuple or None."""
    try:
        fill = shape.fill
        if fill is None:
            return None
        # Solid fill
        if fill.type is not None:
            try:
                fc = fill.fore_color
                if fc and fc.rgb:
                    rgb = fc.rgb
                    return (rgb[0], rgb[1], rgb[2])
            except (AttributeError, TypeError, ValueError):
                pass
    except Exception:
        pass

    # Try XML-level extraction
    try:
        sp_elem = shape._element
        solid = sp_elem.find('.//' + qn('a:solidFill'))
        if solid is not None:
            srgb = solid.find(qn('a:srgbClr'))
            if srgb is not None:
                val = srgb.get('val')
                if val and len(val) == 6:
                    return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
    except Exception:
        pass

    return None


def extract_font_info(run):
    """Extract font name and size from a text run."""
    font_name = None
    font_size_pt = None
    try:
        if run.font:
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size_pt = round(run.font.size / 12700)  # EMU to pt
    except Exception:
        pass
    return font_name or 'Calibri', font_size_pt or 9


def get_pil_font(font_name, font_size_px):
    """Try to load a TrueType font, falling back to default."""
    if not HAS_PIL:
        return None
    try:
        return ImageFont.truetype(font_name, font_size_px)
    except (OSError, IOError):
        pass
    # Try common macOS paths
    mac_fonts = {
        'Georgia': '/System/Library/Fonts/Supplemental/Georgia.ttf',
        'Calibri': '/System/Library/Fonts/Supplemental/Calibri.ttf',
        'Courier New': '/System/Library/Fonts/Supplemental/Courier New.ttf',
        'Arial': '/System/Library/Fonts/Supplemental/Arial.ttf',
    }
    path = mac_fonts.get(font_name)
    if path and os.path.exists(path):
        try:
            return ImageFont.truetype(path, font_size_px)
        except (OSError, IOError):
            pass
    try:
        return ImageFont.load_default()
    except Exception:
        return None


def write_yaml_manual(data, path):
    """Write YAML-like output without requiring PyYAML."""
    if HAS_YAML:
        with open(path, 'w') as f:
            yaml.dump(data, f, default_flow_style=False, sort_keys=False, allow_unicode=True)
        return

    # Manual YAML serializer for simple structures
    def _serialize(obj, indent=0):
        prefix = '  ' * indent
        lines = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                if isinstance(v, (dict, list)):
                    lines.append(f"{prefix}{k}:")
                    lines.extend(_serialize(v, indent + 1))
                else:
                    lines.append(f"{prefix}{k}: {_format_val(v)}")
        elif isinstance(obj, list):
            if not obj:
                return [f"{prefix}[]"]
            for item in obj:
                if isinstance(item, dict):
                    first = True
                    for k, v in item.items():
                        if first:
                            if isinstance(v, (dict, list)):
                                lines.append(f"{prefix}- {k}:")
                                lines.extend(_serialize(v, indent + 2))
                            else:
                                lines.append(f"{prefix}- {k}: {_format_val(v)}")
                            first = False
                        else:
                            if isinstance(v, (dict, list)):
                                lines.append(f"{prefix}  {k}:")
                                lines.extend(_serialize(v, indent + 2))
                            else:
                                lines.append(f"{prefix}  {k}: {_format_val(v)}")
                else:
                    lines.append(f"{prefix}- {_format_val(item)}")
        return lines

    def _format_val(v):
        if v is None:
            return 'null'
        if isinstance(v, bool):
            return 'true' if v else 'false'
        if isinstance(v, (int, float)):
            return str(v)
        s = str(v)
        if ':' in s or '#' in s or s.startswith('{') or s.startswith('['):
            return f'"{s}"'
        return s

    with open(path, 'w') as f:
        f.write('\n'.join(_serialize(data)) + '\n')


# ===========================================================================
# PPTX Renderer
# ===========================================================================

def render_pptx(input_path, output_dir, dpi):
    """Render each slide of a PPTX as a PNG with shape-level analysis."""
    if not HAS_PPTX:
        print("ERROR: python-pptx required for PPTX rendering. Install: pip install python-pptx", file=sys.stderr)
        sys.exit(1)
    if not HAS_PIL:
        print("ERROR: Pillow required for PPTX rendering. Install: pip install Pillow", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(input_path))
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    img_w = emu_to_px(slide_w_emu, dpi)
    img_h = emu_to_px(slide_h_emu, dpi)

    defects = []
    slide_count = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        img = Image.new('RGB', (img_w, img_h), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Draw slide boundary (thin gray)
        draw.rectangle(
            [0, 0, img_w - 1, img_h - 1],
            outline=(200, 200, 200),
            width=1,
        )

        # Collect shape bounding boxes for overlap detection
        shape_boxes = []  # list of (name, left, top, width, height) in EMU

        for shape in slide.shapes:
            s_left = shape.left or 0
            s_top = shape.top or 0
            s_width = shape.width or 0
            s_height = shape.height or 0

            shape_boxes.append((
                shape.name or f"Shape_{shape.shape_id}",
                s_left, s_top, s_width, s_height,
            ))

            px_left = emu_to_px(s_left, dpi)
            px_top = emu_to_px(s_top, dpi)
            px_right = emu_to_px(s_left + s_width, dpi)
            px_bottom = emu_to_px(s_top + s_height, dpi)

            # --- Picture shapes ---
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                draw.rectangle(
                    [px_left, px_top, px_right, px_bottom],
                    fill=(220, 220, 220),
                    outline=(180, 180, 180),
                    width=1,
                )
                # Draw [IMG] label centered
                label = "[IMG]"
                font = get_pil_font('Arial', max(10, (px_bottom - px_top) // 4))
                try:
                    bbox = draw.textbbox((0, 0), label, font=font)
                    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
                except (AttributeError, TypeError):
                    tw, th = 30, 12
                cx = px_left + (px_right - px_left - tw) // 2
                cy = px_top + (px_bottom - px_top - th) // 2
                draw.text((cx, cy), label, fill=(120, 120, 120), font=font)
                continue

            # --- Table shapes ---
            if shape.has_table:
                _render_table(draw, shape, dpi, px_left, px_top, px_right, px_bottom)
                continue

            # --- General shapes (text boxes, rounded rectangles, etc.) ---
            fill_color = extract_fill_color(shape)
            fill_rgb = fill_color if fill_color else None

            # Draw bounding box
            draw.rectangle(
                [px_left, px_top, px_right, px_bottom],
                fill=fill_rgb,
                outline=(180, 180, 180),
                width=1,
            )

            # Boundary violation detection
            if (s_left + s_width > slide_w_emu) or (s_top + s_height > slide_h_emu) or s_left < 0 or s_top < 0:
                defects.append({
                    'page': slide_idx,
                    'category': 'boundary',
                    'severity': 'ERROR',
                    'shapes': [shape.name or f"Shape_{shape.shape_id}"],
                    'description': 'Shape extends beyond slide boundary',
                    'coordinates': {
                        'x': s_left, 'y': s_top,
                        'w': s_width, 'h': s_height,
                    },
                })
                # Draw red border for boundary violation
                draw.rectangle(
                    [px_left, px_top, px_right, px_bottom],
                    outline=(255, 0, 0),
                    width=2,
                )

            # --- Text rendering with word-wrap and overflow detection ---
            if shape.has_text_frame:
                overflow = _render_text_frame(
                    draw, shape, dpi,
                    s_left, s_top, s_width, s_height,
                    px_left, px_top, px_right, px_bottom,
                )
                if overflow:
                    defects.append({
                        'page': slide_idx,
                        'category': 'overflow',
                        'severity': 'ERROR',
                        'shapes': [shape.name or f"Shape_{shape.shape_id}"],
                        'description': f"Text extends {overflow} EMU below container",
                        'coordinates': {
                            'x': s_left, 'y': s_top,
                            'w': s_width, 'h': s_height,
                        },
                    })

        # --- Overlap detection ---
        for i in range(len(shape_boxes)):
            for j in range(i + 1, len(shape_boxes)):
                name_i, l_i, t_i, w_i, h_i = shape_boxes[i]
                name_j, l_j, t_j, w_j, h_j = shape_boxes[j]
                if shapes_overlap((l_i, t_i, w_i, h_i), (l_j, t_j, w_j, h_j)):
                    defects.append({
                        'page': slide_idx,
                        'category': 'overlap',
                        'severity': 'WARNING',
                        'shapes': [name_i, name_j],
                        'description': f"Shapes '{name_i}' and '{name_j}' overlap",
                        'coordinates': {
                            'x': max(l_i, l_j),
                            'y': max(t_i, t_j),
                            'w': min(l_i + w_i, l_j + w_j) - max(l_i, l_j),
                            'h': min(t_i + h_i, t_j + h_j) - max(t_i, t_j),
                        },
                    })
                    # Draw yellow border on overlapping shapes
                    for name, sl, st, sw, sh in [shape_boxes[i], shape_boxes[j]]:
                        draw.rectangle(
                            [emu_to_px(sl, dpi), emu_to_px(st, dpi),
                             emu_to_px(sl + sw, dpi), emu_to_px(st + sh, dpi)],
                            outline=(255, 255, 0),
                            width=2,
                        )

        # Save slide image
        out_path = output_dir / f"slide_{slide_idx:02d}.png"
        img.save(str(out_path), 'PNG')
        print(f"  Rendered: {out_path.name} ({img_w}x{img_h})")

    return slide_count, defects


def _render_table(draw, shape, dpi, px_left, px_top, px_right, px_bottom):
    """Render a table shape with header row, alternating fills, and cell text."""
    table = shape.table
    num_rows = len(table.rows)
    num_cols = len(table.columns)

    if num_rows == 0 or num_cols == 0:
        return

    # Calculate row heights and column widths in pixels
    total_h = px_bottom - px_top
    total_w = px_right - px_left

    # Use table's row/col dimensions if available
    row_heights_emu = []
    for row in table.rows:
        try:
            row_heights_emu.append(row.height or 0)
        except Exception:
            row_heights_emu.append(0)

    col_widths_emu = []
    for col in table.columns:
        try:
            col_widths_emu.append(col.width or 0)
        except Exception:
            col_widths_emu.append(0)

    # If dimensions are zero, distribute evenly
    total_row_emu = sum(row_heights_emu)
    total_col_emu = sum(col_widths_emu)

    if total_row_emu == 0:
        row_heights_px = [total_h // num_rows] * num_rows
    else:
        row_heights_px = [int(h / total_row_emu * total_h) for h in row_heights_emu]

    if total_col_emu == 0:
        col_widths_px = [total_w // num_cols] * num_cols
    else:
        col_widths_px = [int(w / total_col_emu * total_w) for w in col_widths_emu]

    # Draw cells
    y_offset = px_top
    for row_idx, row in enumerate(table.rows):
        x_offset = px_left
        rh = row_heights_px[row_idx] if row_idx < len(row_heights_px) else (total_h // num_rows)

        for col_idx, cell in enumerate(row.cells):
            cw = col_widths_px[col_idx] if col_idx < len(col_widths_px) else (total_w // num_cols)

            # Fill color
            if row_idx == 0:
                fill = (60, 60, 80)  # Dark header
                text_color = (255, 255, 255)
            elif row_idx % 2 == 0:
                fill = (240, 240, 245)
                text_color = (40, 40, 40)
            else:
                fill = (255, 255, 255)
                text_color = (40, 40, 40)

            draw.rectangle(
                [x_offset, y_offset, x_offset + cw, y_offset + rh],
                fill=fill,
                outline=(180, 180, 180),
                width=1,
            )

            # Cell text
            cell_text = cell.text.strip() if cell.text else ""
            if cell_text:
                font_size_px = max(8, rh // 3)
                font = get_pil_font('Calibri', font_size_px)

                # Wrap text to fit cell
                chars_per_line = max(1, cw // max(1, font_size_px // 2))
                wrapped = textwrap.wrap(cell_text, width=chars_per_line)
                text_y = y_offset + 2
                for line in wrapped[:3]:  # Cap at 3 lines per cell
                    if text_y + font_size_px > y_offset + rh:
                        break
                    draw.text((x_offset + 3, text_y), line, fill=text_color, font=font)
                    text_y += font_size_px + 1

            x_offset += cw
        y_offset += rh


def _render_text_frame(draw, shape, dpi, s_left, s_top, s_width, s_height,
                       px_left, px_top, px_right, px_bottom):
    """
    Render text from shape.text_frame with word-wrapping.
    Returns overflow amount in EMU if text overflows, else None.
    """
    tf = shape.text_frame
    y_cursor_emu = s_top  # Track vertical position in EMU

    # Account for internal margins
    margin_left_emu = 0
    margin_top_emu = 0
    margin_right_emu = 0
    margin_bottom_emu = 0
    try:
        if tf.margin_left is not None:
            margin_left_emu = tf.margin_left
        if tf.margin_top is not None:
            margin_top_emu = tf.margin_top
        if tf.margin_right is not None:
            margin_right_emu = tf.margin_right
        if tf.margin_bottom is not None:
            margin_bottom_emu = tf.margin_bottom
    except Exception:
        pass

    usable_width_emu = s_width - margin_left_emu - margin_right_emu
    y_cursor_emu = s_top + margin_top_emu
    container_bottom_emu = s_top + s_height - margin_bottom_emu
    overflow_emu = None

    for para in tf.paragraphs:
        para_text = para.text.strip()
        if not para_text:
            # Empty paragraph -- advance by a line
            y_cursor_emu += 120000  # ~small line height
            continue

        # Get font info from the first run (or paragraph-level)
        font_name = 'Calibri'
        font_size_pt = 9
        font_color = (40, 40, 40)
        font_bold = False

        for run in para.runs:
            fn, fs = extract_font_info(run)
            font_name = fn
            font_size_pt = fs
            try:
                if run.font.bold:
                    font_bold = True
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    font_color = (rgb[0], rgb[1], rgb[2])
            except Exception:
                pass
            break  # Use first run's formatting

        char_width_emu = estimate_char_width_emu(font_name, font_size_pt)
        chars_per_line = max(1, int(usable_width_emu / char_width_emu))

        # Line height in EMU (roughly 1.2x font size)
        line_height_emu = int(font_size_pt * 12700 * 1.2)

        # Wrap text
        wrapped_lines = textwrap.wrap(para_text, width=chars_per_line)
        if not wrapped_lines:
            wrapped_lines = [para_text]

        # Render each wrapped line
        font_size_px = max(8, emu_to_px(font_size_pt * 12700, dpi))
        pil_font = get_pil_font(font_name, font_size_px)

        for line in wrapped_lines:
            px_x = emu_to_px(s_left + margin_left_emu, dpi)
            px_y = emu_to_px(y_cursor_emu, dpi)

            # Check if this line overflows the container
            line_is_overflow = (y_cursor_emu + line_height_emu > container_bottom_emu)

            if line_is_overflow and overflow_emu is None:
                overflow_emu = y_cursor_emu + line_height_emu - container_bottom_emu

            # Draw text
            draw.text(
                (px_x, px_y),
                line,
                fill=font_color,
                font=pil_font,
            )

            # If overflowing, draw RED underline
            if line_is_overflow:
                try:
                    bbox = draw.textbbox((px_x, px_y), line, font=pil_font)
                    line_end_x = bbox[2]
                    underline_y = bbox[3] + 1
                except (AttributeError, TypeError):
                    line_end_x = px_x + len(line) * (font_size_px // 2)
                    underline_y = px_y + font_size_px + 1
                draw.line(
                    [(px_x, underline_y), (line_end_x, underline_y)],
                    fill=(255, 0, 0),
                    width=2,
                )

            y_cursor_emu += line_height_emu

        # Paragraph spacing
        y_cursor_emu += int(line_height_emu * 0.3)

    return overflow_emu


# ===========================================================================
# PDF Renderer
# ===========================================================================

def render_pdf(input_path, output_dir, dpi):
    """Render each page of a PDF as a PNG."""
    page_count = 0

    # Strategy 1: pdf2image (poppler)
    if HAS_PDF2IMAGE:
        try:
            images = convert_from_path(str(input_path), dpi=dpi)
            for i, img in enumerate(images, start=1):
                out_path = output_dir / f"page_{i:02d}.png"
                img.save(str(out_path), 'PNG')
                print(f"  Rendered: {out_path.name} ({img.width}x{img.height})")
                page_count += 1
            return page_count, []
        except Exception as e:
            print(f"  pdf2image failed: {e}. Trying fallback...", file=sys.stderr)

    # Strategy 2: pdftoppm (poppler CLI)
    if shutil.which('pdftoppm'):
        try:
            result = subprocess.run(
                [
                    'pdftoppm', '-png', '-r', str(dpi),
                    str(input_path), str(output_dir / 'page'),
                ],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0:
                # pdftoppm outputs page-1.png, page-2.png, etc.
                for png_file in sorted(output_dir.glob('page-*.png')):
                    # Rename to page_01.png format
                    num = png_file.stem.split('-')[-1]
                    new_name = output_dir / f"page_{int(num):02d}.png"
                    png_file.rename(new_name)
                    page_count += 1
                    if HAS_PIL:
                        img = Image.open(str(new_name))
                        print(f"  Rendered: {new_name.name} ({img.width}x{img.height})")
                        img.close()
                    else:
                        print(f"  Rendered: {new_name.name}")
                return page_count, []
            else:
                print(f"  pdftoppm failed: {result.stderr.strip()}", file=sys.stderr)
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"  pdftoppm error: {e}", file=sys.stderr)

    # Strategy 3: sips (macOS native) -- limited, only first page
    if sys.platform == 'darwin' and shutil.which('sips'):
        try:
            out_path = output_dir / "page_01.png"
            result = subprocess.run(
                [
                    'sips', '-s', 'format', 'png',
                    '--resampleHeightWidthMax', str(int(11 * dpi)),
                    str(input_path),
                    '--out', str(out_path),
                ],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0 and out_path.exists():
                print(f"  Rendered: {out_path.name} (sips -- may be first page only)")
                return 1, []
            else:
                print(f"  sips failed: {result.stderr.strip()}", file=sys.stderr)
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"  sips error: {e}", file=sys.stderr)

    print("ERROR: No PDF renderer available.", file=sys.stderr)
    print("  Install one of:", file=sys.stderr)
    print("    pip install pdf2image   (+ brew install poppler)", file=sys.stderr)
    print("    brew install poppler    (for pdftoppm CLI)", file=sys.stderr)
    sys.exit(1)


# ===========================================================================
# DOCX Renderer
# ===========================================================================

def render_docx(input_path, output_dir, dpi):
    """Convert DOCX to PDF, then render via PDF pipeline."""
    temp_pdf = None
    temp_dirs = []  # Track ALL created temp directories for cleanup

    try:
        # Strategy 1: textutil (macOS native)
        if sys.platform == 'darwin' and shutil.which('textutil'):
            temp_dir = tempfile.mkdtemp(prefix='render-doc-audit-')
            temp_dirs.append(temp_dir)
            temp_pdf = Path(temp_dir) / (input_path.stem + '.pdf')
            try:
                result = subprocess.run(
                    [
                        'textutil', '-convert', 'pdf',
                        '-output', str(temp_pdf),
                        str(input_path),
                    ],
                    capture_output=True, text=True, timeout=60,
                )
                if result.returncode == 0 and temp_pdf.exists():
                    print(f"  Converted DOCX to PDF via textutil")
                else:
                    print(f"  textutil failed: {result.stderr.strip()}", file=sys.stderr)
                    temp_pdf = None
            except (subprocess.TimeoutExpired, FileNotFoundError) as e:
                print(f"  textutil error: {e}", file=sys.stderr)
                temp_pdf = None

        # Strategy 2: LibreOffice
        if temp_pdf is None and shutil.which('libreoffice'):
            temp_dir = tempfile.mkdtemp(prefix='render-doc-audit-')
            temp_dirs.append(temp_dir)
            try:
                result = subprocess.run(
                    [
                        'libreoffice', '--headless', '--convert-to', 'pdf',
                        '--outdir', temp_dir,
                        str(input_path),
                    ],
                    capture_output=True, text=True, timeout=120,
                )
                if result.returncode == 0:
                    candidates = list(Path(temp_dir).glob('*.pdf'))
                    if candidates:
                        temp_pdf = candidates[0]
                        print(f"  Converted DOCX to PDF via LibreOffice")
            except (subprocess.TimeoutExpired, FileNotFoundError) as e:
                print(f"  LibreOffice error: {e}", file=sys.stderr)

        if temp_pdf is None or not temp_pdf.exists():
            print("ERROR: Cannot convert DOCX to PDF.", file=sys.stderr)
            print("  Requires: textutil (macOS) or libreoffice", file=sys.stderr)
            sys.exit(1)

        # Render the temporary PDF
        page_count, defects = render_pdf(temp_pdf, output_dir, dpi)
    finally:
        # Clean up ALL temp directories regardless of which strategy succeeded
        for td in temp_dirs:
            shutil.rmtree(td, ignore_errors=True)

    return page_count, defects


# ===========================================================================
# Main
# ===========================================================================

def detect_format(input_path):
    """Auto-detect document format from file extension."""
    ext = input_path.suffix.lower()
    format_map = {
        '.pptx': 'pptx',
        '.pdf': 'pdf',
        '.docx': 'docx',
    }
    return format_map.get(ext)


def main():
    parser = argparse.ArgumentParser(
        description='ACOS Document Screenshot Renderer -- Universal Visual QA Tool.\n'
                    'Renders PPTX, PDF, or DOCX files to PNG screenshots for visual QA.\n'
                    'PPTX rendering includes shape-level analysis: overflow detection,\n'
                    'overlap detection, boundary checks, and text wrap simulation.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            'Examples:\n'
            '  python3 render-doc-audit.py presentation.pptx\n'
            '  python3 render-doc-audit.py report.pdf --dpi 300 --output-dir ./screenshots\n'
            '  python3 render-doc-audit.py document.docx --format docx\n'
            '\n'
            'Dependencies:\n'
            '  PPTX:  python-pptx + Pillow (required)\n'
            '  PDF:   pdf2image + poppler (preferred) | pdftoppm | sips (macOS)\n'
            '  DOCX:  textutil (macOS) or libreoffice (converts to PDF first)\n'
        ),
    )
    parser.add_argument('input_file', help='Path to the document to render')
    parser.add_argument('--output-dir', '-o', default=None,
                        help='Output directory for PNGs (default: {input_dir}/visual-audit/)')
    parser.add_argument('--dpi', type=int, default=150,
                        help='Resolution in DPI (default: 150). US Letter at 150 DPI = '
                             '1275x1650 px, safely under the Anthropic API 2000 px '
                             'many-image cap. Higher DPI risks locking the conversation '
                             'in an unrecoverable image-cap error during downstream Read loops.')
    parser.add_argument('--format', choices=['pptx', 'pdf', 'docx', 'auto'], default='auto',
                        help='Document format (default: auto-detect from extension)')

    args = parser.parse_args()

    # Resolve input file
    input_path = Path(args.input_file).resolve()
    if not input_path.exists():
        print(f"ERROR: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    # Detect format
    doc_format = args.format
    if doc_format == 'auto':
        doc_format = detect_format(input_path)
        if doc_format is None:
            print(f"ERROR: Cannot auto-detect format for '{input_path.name}'.", file=sys.stderr)
            print(f"  Supported extensions: .pptx, .pdf, .docx", file=sys.stderr)
            print(f"  Use --format to specify manually.", file=sys.stderr)
            sys.exit(1)

    # Resolve output directory
    if args.output_dir:
        output_dir = Path(args.output_dir).resolve()
    else:
        output_dir = input_path.parent / 'visual-audit'

    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"ACOS Document Screenshot Renderer")
    print(f"  Input:  {input_path}")
    print(f"  Format: {doc_format.upper()}")
    print(f"  DPI:    {args.dpi}")
    print(f"  Output: {output_dir}")
    print()

    # Dispatch to format-specific renderer
    if doc_format == 'pptx':
        page_count, defects = render_pptx(input_path, output_dir, args.dpi)
        file_prefix = 'slide'
    elif doc_format == 'pdf':
        page_count, defects = render_pdf(input_path, output_dir, args.dpi)
        file_prefix = 'page'
    elif doc_format == 'docx':
        page_count, defects = render_docx(input_path, output_dir, args.dpi)
        file_prefix = 'page'
    else:
        print(f"ERROR: Unsupported format: {doc_format}", file=sys.stderr)
        sys.exit(1)

    # Write audit summary YAML
    summary = {
        'input_file': str(input_path),
        'format': doc_format,
        'pages': page_count,
        'dpi': args.dpi,
        'output_dir': str(output_dir),
        'defects': defects if defects else [],
    }

    summary_path = output_dir / 'audit-summary.yaml'
    write_yaml_manual(summary, summary_path)

    print()
    print(f"Summary: {summary_path}")
    print(f"  Pages rendered: {page_count}")
    if defects:
        print(f"  Defects found:  {len(defects)}")
        for d in defects:
            print(f"    - [{file_prefix} {d['page']}] {d['category']}: {d['description']}")
    else:
        if doc_format == 'pptx':
            print(f"  Defects found:  0 (clean)")
        else:
            print(f"  (Shape-level defect detection only available for PPTX)")

    print()
    print("Done.")

    # Exit with code 1 if any defects were found, matching the contract
    # used by check-pptx-layout.py and check-pdf-layout.py
    if defects:
        sys.exit(1)


if __name__ == '__main__':
    main()
