#!/usr/bin/env python3
"""
ACOS PPTX Layout Checker — Fast Programmatic Pre-Check

Checks a PPTX file for layout defects WITHOUT rendering images.
Pure code-level analysis using python-pptx shape geometry.

Usage:
    python3 check-pptx-layout.py <pptx-file> [-o <report.yaml>] [--strict] [--verbose]

Checks (7):
    1. boundary_violation  — shape extends past slide edge
    2. shape_overlap       — bounding box intersection between shapes
    3. text_overflow        — estimated text content exceeds container
    4. footer_collision    — content shape bottom near footer divider
    5. margin_compliance   — content shapes outside safe margins
    6. cross_slide_consistency — structural elements drift across slides
    7. callout_consistency — child text boxes escape parent rounded rect

Dependencies: python-pptx, PyYAML (stdlib otherwise)
"""
import argparse
import math
import re
import sys
from collections import defaultdict
from pathlib import Path

try:
    import yaml
except ImportError:
    print("ERROR: PyYAML required. Install: pip install PyYAML", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Emu, Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# =============================================================================
# CONSTANTS
# =============================================================================

# Standard 16:9 widescreen dimensions
DEFAULT_SLIDE_WIDTH = 12192000   # 13.333" in EMU
DEFAULT_SLIDE_HEIGHT = 6858000   # 7.5" in EMU

# Footer divider Y position (standard slides)
FOOTER_DIVIDER_Y = 4709160

# Footer collision threshold
FOOTER_COLLISION_THRESHOLD = 50000  # EMU

# Margins (content safe zone)
MARGIN_LEFT = 457200     # 0.5" in EMU
MARGIN_RIGHT_LIMIT = 8686800   # 9.5" in EMU (for standard 10" slides)
# NOTE: For widescreen (13.333"), right limit = slide_width - 457200

# Overlap thresholds
SMALL_OVERLAP_THRESHOLD = 50000  # EMU — below this is WARNING, above is ERROR

# Cross-slide consistency drift threshold
CONSISTENCY_DRIFT_THRESHOLD = 50000  # EMU

# Character width estimation table (EMU per character at given font size)
# Formula: font_size_pt * multiplier * 12700 (EMU per point)
# Multipliers by font family category:
#   Monospace (Courier New, Consolas): 0.62
#   Proportional serif (Georgia, Times): 0.50
#   Proportional sans (Calibri, Arial, Helvetica): 0.48
#   Narrow (Arial Narrow, Calibri Light): 0.40
#   Wide (Impact, Verdana): 0.55
CHAR_WIDTH_MULTIPLIERS = {
    # Monospace
    "Courier New": 0.62,
    "Consolas": 0.62,
    "Courier": 0.62,
    "Monaco": 0.62,
    "Menlo": 0.62,
    # Serif proportional
    "Georgia": 0.50,
    "Times New Roman": 0.48,
    "Garamond": 0.45,
    "Palatino": 0.50,
    # Sans proportional
    "Calibri": 0.48,
    "Arial": 0.50,
    "Helvetica": 0.50,
    "Segoe UI": 0.48,
    "Verdana": 0.55,
    "Tahoma": 0.50,
    # Narrow
    "Arial Narrow": 0.40,
}
DEFAULT_CHAR_WIDTH_MULTIPLIER = 0.48

# Line height multiplier (relative to font size)
LINE_HEIGHT_MULTIPLIER = 1.4  # 140% of font size is standard single spacing


# =============================================================================
# GEOMETRY HELPERS
# =============================================================================

def get_bbox(shape):
    """Return (left, top, right, bottom) bounding box for a shape.

    Returns None if any position property is None (placeholder shapes
    with inherited positions from slide layout).
    """
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    return (shape.left, shape.top,
            shape.left + shape.width, shape.top + shape.height)


def bbox_intersects(a, b):
    """Check if two bounding boxes (l, t, r, b) intersect. Returns overlap area in EMU^2 or 0."""
    x_overlap = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    y_overlap = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    return x_overlap * y_overlap


def bbox_contains(outer, inner):
    """Check if outer bbox fully contains inner bbox."""
    return (outer[0] <= inner[0] and outer[1] <= inner[1] and
            outer[2] >= inner[2] and outer[3] >= inner[3])


def is_full_slide_rect(shape, slide_width, slide_height):
    """Check if shape is a background fill (covers most of the slide)."""
    coverage_w = shape.width / slide_width if slide_width > 0 else 0
    coverage_h = shape.height / slide_height if slide_height > 0 else 0
    return coverage_w > 0.9 and coverage_h > 0.9


def overlap_linear_extent(a, b):
    """Return the max linear overlap (in EMU) between two bboxes along either axis."""
    x_overlap = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    y_overlap = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    return max(x_overlap, y_overlap)


def char_width_emu(font_name, font_size_pt):
    """Estimate character width in EMU for a given font and size."""
    multiplier = CHAR_WIDTH_MULTIPLIERS.get(font_name, DEFAULT_CHAR_WIDTH_MULTIPLIER)
    return int(font_size_pt * multiplier * 12700)


def line_height_emu(font_size_pt):
    """Estimate line height in EMU for a given font size."""
    return int(font_size_pt * LINE_HEIGHT_MULTIPLIER * 12700)


# =============================================================================
# TEXT OVERFLOW ESTIMATION
# =============================================================================

def estimate_text_lines(text, available_width_emu, font_name, font_size_pt):
    """Estimate number of visual lines needed for text with word wrapping.

    Returns (total_lines, chars_per_line).
    """
    if not text or available_width_emu <= 0 or font_size_pt <= 0:
        return (0, 0)

    cw = char_width_emu(font_name, font_size_pt)
    if cw <= 0:
        return (0, 0)

    chars_per_line = max(1, int(available_width_emu / cw))

    # Process each paragraph (newline-separated) independently
    total_lines = 0
    paragraphs = text.split("\n")
    for para in paragraphs:
        para = para.strip()
        if not para:
            total_lines += 1  # blank line
            continue

        # Word-wrap estimation
        words = para.split()
        current_line_chars = 0
        para_lines = 1

        for word in words:
            word_len = len(word)
            if current_line_chars == 0:
                # First word on the line
                current_line_chars = word_len
            elif current_line_chars + 1 + word_len <= chars_per_line:
                # Word fits on current line (with space)
                current_line_chars += 1 + word_len
            else:
                # Wrap to next line
                para_lines += 1
                current_line_chars = word_len

            # Handle words longer than a full line
            if word_len > chars_per_line:
                extra_lines = (word_len - 1) // chars_per_line
                para_lines += extra_lines

        total_lines += para_lines

    return (total_lines, chars_per_line)


# =============================================================================
# LAYOUT CHECKER
# =============================================================================

class LayoutChecker:
    def __init__(self, pptx_path, strict=False, verbose=False):
        self.pptx_path = str(pptx_path)
        self.strict = strict
        self.verbose = verbose
        self.defects = []
        self.total_shapes = 0
        self.checks_run = 7

        try:
            self.prs = Presentation(pptx_path)
        except Exception as e:
            self.prs = None
            self.defects.append({
                "slide": 0,
                "check": "file_error",
                "severity": "ERROR",
                "shapes": [],
                "message": "Cannot open PPTX: {}".format(e),
                "fix_suggestion": "Verify the file is a valid .pptx",
            })
            return

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        # Compute right margin limit based on actual slide width
        self.margin_right_limit = self.slide_width - MARGIN_LEFT

    def add_defect(self, slide_num, check, severity, shapes, message, fix_suggestion=""):
        self.defects.append({
            "slide": slide_num,
            "check": check,
            "severity": severity,
            "shapes": shapes if isinstance(shapes, list) else [shapes],
            "message": message,
            "fix_suggestion": fix_suggestion,
        })

    def _get_dominant_font(self, shape):
        """Extract the most common font name and largest font size from a text shape."""
        if not shape.has_text_frame:
            return (None, 12)

        font_sizes = []
        font_names = []

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(run.font.size.pt)
                if run.font.name:
                    font_names.append(run.font.name)

        font_name = font_names[0] if font_names else "Calibri"
        font_size = max(font_sizes) if font_sizes else 12

        return (font_name, font_size)

    def _get_text_margins(self, shape):
        """Get internal text margins from shape text_frame."""
        if not shape.has_text_frame:
            return (0, 0, 0, 0)
        tf = shape.text_frame
        return (
            tf.margin_left if tf.margin_left is not None else 91440,
            tf.margin_top if tf.margin_top is not None else 45720,
            tf.margin_right if tf.margin_right is not None else 91440,
            tf.margin_bottom if tf.margin_bottom is not None else 45720,
        )

    # ----- Check 1: Boundary Violations -----

    def check_boundary_violations(self):
        """Any shape where left+width > slide_width or top+height > slide_height."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                self.total_shapes += 1
                # Skip placeholder shapes with inherited positions (None values)
                if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                    continue
                right = shape.left + shape.width
                bottom = shape.top + shape.height

                if shape.left < 0:
                    overshoot = abs(shape.left)
                    self.add_defect(
                        slide_idx, "boundary_violation", "ERROR",
                        [shape.name],
                        "Shape extends {} EMU past left edge (L={})".format(overshoot, shape.left),
                        "Move shape right by {} EMU".format(overshoot),
                    )
                if shape.top < 0:
                    overshoot = abs(shape.top)
                    self.add_defect(
                        slide_idx, "boundary_violation", "ERROR",
                        [shape.name],
                        "Shape extends {} EMU past top edge (T={})".format(overshoot, shape.top),
                        "Move shape down by {} EMU".format(overshoot),
                    )
                if right > self.slide_width:
                    overshoot = right - self.slide_width
                    self.add_defect(
                        slide_idx, "boundary_violation", "ERROR",
                        [shape.name],
                        "Shape extends {} EMU past right edge (R={}, max={})".format(
                            overshoot, right, self.slide_width),
                        "Reduce width by {} EMU".format(overshoot),
                    )
                if bottom > self.slide_height:
                    overshoot = bottom - self.slide_height
                    self.add_defect(
                        slide_idx, "boundary_violation", "ERROR",
                        [shape.name],
                        "Shape extends {} EMU past bottom edge (B={}, max={})".format(
                            overshoot, bottom, self.slide_height),
                        "Reduce height by {} EMU".format(overshoot),
                    )

    # ----- Check 2: Shape-to-Shape Overlap -----

    def check_shape_overlap(self):
        """For each pair of shapes on the same slide, check bounding box intersection."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            shapes = list(slide.shapes)
            # Skip shapes with None positions (placeholder shapes with inherited layout)
            bboxes = [(s, get_bbox(s)) for s in shapes]
            bboxes = [(s, bb) for s, bb in bboxes if bb is not None]

            for i in range(len(bboxes)):
                for j in range(i + 1, len(bboxes)):
                    s_a, bb_a = bboxes[i]
                    s_b, bb_b = bboxes[j]

                    # Skip if one fully contains the other (parent-child)
                    if bbox_contains(bb_a, bb_b) or bbox_contains(bb_b, bb_a):
                        continue

                    # Skip if either is a full-slide background
                    if (is_full_slide_rect(s_a, self.slide_width, self.slide_height) or
                            is_full_slide_rect(s_b, self.slide_width, self.slide_height)):
                        continue

                    area = bbox_intersects(bb_a, bb_b)
                    if area > 0:
                        linear = overlap_linear_extent(bb_a, bb_b)
                        if linear < SMALL_OVERLAP_THRESHOLD:
                            severity = "WARNING"
                        else:
                            severity = "ERROR"

                        self.add_defect(
                            slide_idx, "shape_overlap", severity,
                            [s_a.name, s_b.name],
                            "Shapes overlap by {} EMU (area={} EMU^2)".format(linear, area),
                            "Separate shapes by at least {} EMU".format(linear),
                        )

    # ----- Check 3: Text Overflow Estimation -----

    def check_text_overflow(self):
        """Estimate whether text content fits in its container shape."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # Skip placeholder shapes with inherited positions (None values)
                if shape.width is None or shape.height is None:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                font_name, font_size = self._get_dominant_font(shape)
                ml, mt, mr, mb = self._get_text_margins(shape)

                available_width = shape.width - ml - mr
                available_height = shape.height - mt - mb

                if available_width <= 0 or available_height <= 0:
                    continue

                num_lines, chars_per_line = estimate_text_lines(
                    text, available_width, font_name, font_size
                )

                if num_lines == 0:
                    continue

                required_height = num_lines * line_height_emu(font_size)

                if required_height > available_height:
                    # Use a tolerance of 10% before flagging
                    if required_height > available_height * 1.1:
                        fits_lines = max(1, int(available_height / line_height_emu(font_size)))
                        self.add_defect(
                            slide_idx, "text_overflow", "ERROR",
                            [shape.name],
                            "Text estimated at {} lines but container fits {} lines "
                            "(H={}, needs ~{})".format(
                                num_lines, fits_lines, available_height, required_height),
                            "Increase height to {} EMU".format(int(required_height * 1.1)),
                        )

    # ----- Check 4: Footer Collision -----

    def check_footer_collision(self):
        """Content shape whose bottom edge is within threshold of footer divider."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                # Skip placeholder shapes with inherited positions (None values)
                if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                    continue
                # Skip shapes that ARE the footer (below divider line)
                if shape.top >= FOOTER_DIVIDER_Y:
                    continue
                # Skip full-slide backgrounds
                if is_full_slide_rect(shape, self.slide_width, self.slide_height):
                    continue

                bottom = shape.top + shape.height
                distance = FOOTER_DIVIDER_Y - bottom

                if 0 < distance < FOOTER_COLLISION_THRESHOLD:
                    self.add_defect(
                        slide_idx, "footer_collision", "WARNING",
                        [shape.name],
                        "Shape bottom within {} EMU of footer divider "
                        "(B={}, footer={})".format(distance, bottom, FOOTER_DIVIDER_Y),
                        "Move shape up by {} EMU or reduce height".format(
                            FOOTER_COLLISION_THRESHOLD - distance),
                    )
                elif distance < 0 and bottom > FOOTER_DIVIDER_Y:
                    # Shape crosses the footer line but starts above it
                    overlap_amount = bottom - FOOTER_DIVIDER_Y
                    if overlap_amount < shape.height * 0.5:
                        # Not a footer element itself — it's overlapping
                        self.add_defect(
                            slide_idx, "footer_collision", "WARNING",
                            [shape.name],
                            "Shape crosses footer divider by {} EMU "
                            "(B={}, footer={})".format(
                                overlap_amount, bottom, FOOTER_DIVIDER_Y),
                            "Reduce height by {} EMU".format(overlap_amount),
                        )

    # ----- Check 5: Margin Compliance -----

    def check_margin_compliance(self):
        """Check that content shapes respect left/right margins."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                # Skip placeholder shapes with inherited positions (None values)
                if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                    continue
                # Skip full-slide backgrounds and decorative elements
                if is_full_slide_rect(shape, self.slide_width, self.slide_height):
                    continue

                # Only check shapes that have text (content shapes)
                if not shape.has_text_frame:
                    continue
                if not shape.text_frame.text.strip():
                    continue

                # Left margin check
                if shape.left < MARGIN_LEFT:
                    violation = MARGIN_LEFT - shape.left
                    self.add_defect(
                        slide_idx, "margin_compliance", "WARNING",
                        [shape.name],
                        "Shape starts {} EMU inside left margin "
                        "(L={}, min={})".format(violation, shape.left, MARGIN_LEFT),
                        "Move shape right by {} EMU".format(violation),
                    )

                # Right margin check
                right = shape.left + shape.width
                if right > self.margin_right_limit:
                    violation = right - self.margin_right_limit
                    self.add_defect(
                        slide_idx, "margin_compliance", "WARNING",
                        [shape.name],
                        "Shape extends {} EMU past right margin "
                        "(R={}, max={})".format(
                            violation, right, self.margin_right_limit),
                        "Reduce width by {} EMU".format(violation),
                    )

    # ----- Check 6: Cross-Slide Consistency -----

    def check_cross_slide_consistency(self):
        """Compare structural element positions across slides for drift."""
        # Collect structural elements per slide
        header_positions = {}   # slide_idx -> (left, top, width, height)
        footer_positions = {}

        for slide_idx, slide in enumerate(self.prs.slides, 1):
            # Skip first slide (title slide often has different layout)
            if slide_idx == 1:
                continue

            for shape in slide.shapes:
                # Skip placeholder shapes with inherited positions (None values)
                if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                    continue
                name_lower = shape.name.lower() if shape.name else ""

                # Detect header bars (typically at top, full width or near-full).
                # Require a real header NAME signal — keying off "header in the
                # shape name" rather than "lacks a text frame" prevents unrelated
                # full-width graphics (pictures, lines, decorative bars) from being
                # mistaken for header bars and emitting spurious drift warnings.
                if shape.top < 500000 and shape.width > self.slide_width * 0.5:
                    if "header" in name_lower:
                        header_positions[slide_idx] = (
                            shape.left, shape.top, shape.width, shape.height
                        )

                # Detect footer elements (at bottom of slide)
                if shape.top > FOOTER_DIVIDER_Y:
                    if not footer_positions.get(slide_idx):
                        footer_positions[slide_idx] = (
                            shape.left, shape.top, shape.width, shape.height
                        )

        # Check header consistency
        if len(header_positions) >= 2:
            positions = list(header_positions.values())
            ref = positions[0]
            for slide_idx, pos in header_positions.items():
                for dim_idx, dim_name in enumerate(["left", "top", "width", "height"]):
                    drift = abs(pos[dim_idx] - ref[dim_idx])
                    if drift > CONSISTENCY_DRIFT_THRESHOLD:
                        self.add_defect(
                            slide_idx, "cross_slide_consistency", "WARNING",
                            ["header_bar"],
                            "Header bar {} drifts {} EMU from reference "
                            "(this={}, ref={})".format(
                                dim_name, drift, pos[dim_idx], ref[dim_idx]),
                            "Align header bar {} to {} EMU".format(dim_name, ref[dim_idx]),
                        )

        # Check footer consistency
        if len(footer_positions) >= 2:
            positions = list(footer_positions.values())
            ref = positions[0]
            for slide_idx, pos in footer_positions.items():
                for dim_idx, dim_name in enumerate(["left", "top", "width", "height"]):
                    drift = abs(pos[dim_idx] - ref[dim_idx])
                    if drift > CONSISTENCY_DRIFT_THRESHOLD:
                        self.add_defect(
                            slide_idx, "cross_slide_consistency", "WARNING",
                            ["footer"],
                            "Footer {} drifts {} EMU from reference "
                            "(this={}, ref={})".format(
                                dim_name, drift, pos[dim_idx], ref[dim_idx]),
                            "Align footer {} to {} EMU".format(dim_name, ref[dim_idx]),
                        )

    # ----- Check 7: Callout Internal Consistency -----

    def check_callout_consistency(self):
        """For rounded rectangles containing text boxes, check children are within parent."""
        for slide_idx, slide in enumerate(self.prs.slides, 1):
            shapes = list(slide.shapes)

            # Find rounded rectangles (potential callout parents)
            rounded_rects = []
            for s in shapes:
                try:
                    if (hasattr(s, 'shape_type') and
                            s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE):
                        # Check if it's a rounded rectangle via XML
                        sp_elem = s._element
                        prstGeom = sp_elem.find(
                            './/{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                        if prstGeom is not None:
                            prst = prstGeom.get('prst', '')
                            if 'roundRect' in prst or 'snipRndRect' in prst:
                                rounded_rects.append(s)
                except Exception:
                    continue

            if not rounded_rects:
                continue

            # For each rounded rect, find text boxes that appear to be children
            for parent in rounded_rects:
                parent_bb = get_bbox(parent)
                if parent_bb is None:
                    continue

                for child in shapes:
                    if child is parent:
                        continue
                    if not child.has_text_frame:
                        continue
                    if not child.text_frame.text.strip():
                        continue

                    child_bb = get_bbox(child)
                    if child_bb is None:
                        continue

                    # Check if child overlaps significantly with parent
                    area = bbox_intersects(parent_bb, child_bb)
                    if area <= 0:
                        continue

                    # Calculate what fraction of the child is inside the parent
                    child_area = child.width * child.height
                    if child_area <= 0:
                        continue
                    inside_fraction = area / child_area

                    # If child is mostly inside (>50%) but not fully contained
                    if inside_fraction > 0.5 and not bbox_contains(parent_bb, child_bb):
                        # Calculate how much it escapes
                        escape_left = max(0, parent_bb[0] - child_bb[0])
                        escape_top = max(0, parent_bb[1] - child_bb[1])
                        escape_right = max(0, child_bb[2] - parent_bb[2])
                        escape_bottom = max(0, child_bb[3] - parent_bb[3])
                        max_escape = max(escape_left, escape_top, escape_right, escape_bottom)

                        self.add_defect(
                            slide_idx, "callout_consistency", "ERROR",
                            [child.name, parent.name],
                            "Text box '{}' escapes callout '{}' "
                            "by {} EMU ({:.0%} inside)".format(
                                child.name, parent.name, max_escape, inside_fraction),
                            "Move/resize '{}' to fit within '{}'".format(
                                child.name, parent.name),
                        )

    # ----- Run All Checks -----

    def run_all(self):
        """Execute all 7 layout checks."""
        if self.prs is None:
            return

        self.check_boundary_violations()
        self.check_shape_overlap()
        self.check_text_overflow()
        self.check_footer_collision()
        self.check_margin_compliance()
        self.check_cross_slide_consistency()
        self.check_callout_consistency()

    def build_report(self):
        """Build the YAML-serializable report dict."""
        errors = [d for d in self.defects if d["severity"] == "ERROR"]
        warnings = [d for d in self.defects if d["severity"] == "WARNING"]

        # Count shapes that have no issues
        defect_shape_set = set()
        for d in self.defects:
            for s in d["shapes"]:
                defect_shape_set.add((d["slide"], s))

        passes = max(0, self.total_shapes - len(defect_shape_set))

        report = {
            "file": self.pptx_path,
            "slide_count": len(self.prs.slides) if self.prs else 0,
            "shape_count": self.total_shapes,
            "checks_run": self.checks_run,
            "defect_count": len(self.defects),
            "defects": self.defects,
            "summary": {
                "errors": len(errors),
                "warnings": len(warnings),
                "passes": passes,
            },
        }
        return report

    def print_human_summary(self, report):
        """Print human-readable summary to stdout."""
        filename = Path(self.pptx_path).name
        print("")
        print("PPTX Layout Check: {}".format(filename))
        print("=" * 63)
        print("{} slides, {} shapes checked\n".format(
            report['slide_count'], report['shape_count']))

        errors = [d for d in self.defects if d["severity"] == "ERROR"]
        warnings = [d for d in self.defects if d["severity"] == "WARNING"]

        if errors:
            print("ERRORS ({}):".format(len(errors)))
            for d in errors:
                shapes_str = ", ".join(d["shapes"])
                print("  Slide {:>2}:  {} -- {}".format(
                    d["slide"], shapes_str, d["message"]))
            print()

        if warnings:
            print("WARNINGS ({}):".format(len(warnings)))
            for d in warnings:
                shapes_str = ", ".join(d["shapes"])
                print("  Slide {:>2}:  {} -- {}".format(
                    d["slide"], shapes_str, d["message"]))
            print()

        if not errors and not warnings:
            print("No defects found.\n")

        error_count = report["summary"]["errors"]
        warning_count = report["summary"]["warnings"]

        if error_count > 0:
            result = "FAIL"
        elif warning_count > 0 and self.strict:
            result = "FAIL"
        elif warning_count > 0:
            result = "WARN"
        else:
            result = "PASS"

        print("RESULT: {} ({} errors, {} warnings)".format(
            result, error_count, warning_count))
        return result


# =============================================================================
# MAIN
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="ACOS PPTX Layout Checker -- fast programmatic pre-check for layout defects",
        epilog="Checks: boundary violations, shape overlap, text overflow, footer collision, "
               "margin compliance, cross-slide consistency, callout internal consistency.",
    )
    parser.add_argument("pptx_file", help="Path to the .pptx file to check")
    parser.add_argument("-o", "--output", metavar="REPORT",
                        help="Write YAML report to file (default: stdout summary only)")
    parser.add_argument("--strict", action="store_true",
                        help="Treat warnings as failures (exit code 1)")
    parser.add_argument("--verbose", action="store_true",
                        help="Print detailed per-shape analysis to stderr")
    args = parser.parse_args()

    pptx_path = Path(args.pptx_file)
    if not pptx_path.exists():
        print("ERROR: File not found: {}".format(pptx_path), file=sys.stderr)
        sys.exit(2)
    if not pptx_path.suffix.lower() == ".pptx":
        print("WARNING: File does not have .pptx extension: {}".format(pptx_path),
              file=sys.stderr)

    checker = LayoutChecker(pptx_path, strict=args.strict, verbose=args.verbose)
    checker.run_all()
    report = checker.build_report()

    # Write YAML report if requested
    if args.output:
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        with open(args.output, "w") as f:
            yaml.dump(report, f, default_flow_style=False, sort_keys=False,
                      allow_unicode=True)
        print("Report written to: {}".format(args.output), file=sys.stderr)

    # Always print human summary to stdout
    result = checker.print_human_summary(report)

    # Verbose: dump full YAML to stderr
    if args.verbose and not args.output:
        print("\n--- Full Report (YAML) ---", file=sys.stderr)
        yaml.dump(report, sys.stderr, default_flow_style=False, sort_keys=False,
                  allow_unicode=True)

    # Exit code
    if result == "FAIL":
        sys.exit(1)
    else:
        sys.exit(0)


if __name__ == "__main__":
    main()
