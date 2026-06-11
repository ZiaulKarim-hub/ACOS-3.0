#!/usr/bin/env python3
"""
pptx-cleanup.py — post-build PPTX hygiene + optional inline generator.

Modes:
  --pptx <file>                     Strip theme shadows, enforce explicit backgrounds on
                                    text shapes, attempt Cormorant Garamond embedding.
  --inline-generate ...             Full fallback generator used when data-to-pptx.py
                                    rejects our schema. Builds PPTX directly from the
                                    content YAML + design-spec using python-pptx.

Usage:
    pptx-cleanup.py --pptx <file> [--design-spec <yaml>]
    pptx-cleanup.py --inline-generate --content <yaml> --design-spec <yaml>
                    --template <pptx> --output <pptx>
"""
from __future__ import annotations

import argparse
import shutil
import sys
import zipfile
from pathlib import Path

try:
    import yaml
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from lxml import etree
except ImportError as e:
    sys.stderr.write(f"ERROR: {e.name} required\n")
    sys.exit(1)


THEME_XML_PATH = "ppt/theme/theme1.xml"
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def strip_theme_shadows(pptx_path: Path) -> bool:
    """Remove outerShdw from theme1.xml effectStyleLst; reset effectRef idx=0."""
    if not pptx_path.exists():
        return False
    tmp = pptx_path.with_suffix(".cleaning.pptx")
    shutil.copy(pptx_path, tmp)
    changed = False
    try:
        with zipfile.ZipFile(tmp, "a") as zf:
            try:
                with zf.open(THEME_XML_PATH) as tf:
                    tree = etree.parse(tf)
            except KeyError:
                return False
        root = tree.getroot()
        for shdw in root.findall(".//a:outerShdw", NSMAP):
            parent = shdw.getparent()
            if parent is not None:
                parent.remove(shdw)
                changed = True
        # Reset effectRef idx to 0 where present
        for ref in root.findall(".//a:effectRef", NSMAP):
            if ref.get("idx") and ref.get("idx") != "0":
                ref.set("idx", "0")
                changed = True
        if changed:
            # Rewrite theme1.xml inside the pptx
            with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
                with zipfile.ZipFile(tmp, "r") as zin:
                    for item in zin.namelist():
                        if item == THEME_XML_PATH:
                            zout.writestr(item, etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True))
                        else:
                            zout.writestr(item, zin.read(item))
    finally:
        if tmp.exists():
            tmp.unlink()
    return changed


def enforce_explicit_backgrounds(pptx_path: Path, design_spec: dict | None) -> int:
    """Call fill.solid() on every text-bearing shape whose fill is not explicit."""
    prs = Presentation(pptx_path)
    default_rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if design_spec:
        theme = design_spec.get("theme_colors") or {}
        hex_ = (theme.get("white") or {}).get("rgb_hex", "FFFFFF")
        default_rgb = RGBColor.from_string(hex_.upper())
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                try:
                    f = shape.fill
                    if getattr(f, "type", None) is None:
                        f.solid()
                        f.fore_color.rgb = default_rgb
                        count += 1
                except Exception:
                    pass
    if count:
        prs.save(pptx_path)
    return count


def try_embed_fonts(pptx_path: Path, design_spec: dict | None) -> bool:
    """Best-effort font embedding. On failure, add advisory note to the final slide."""
    # Heavy lift (actual OOXML font embedding) deferred — we surface the fallback advisory.
    if not design_spec:
        return False
    prs = Presentation(pptx_path)
    if not prs.slides:
        return False
    last = prs.slides[list(prs.slides._sldIdLst)[-1].rId] if False else list(prs.slides)[-1]
    from pptx.util import Emu
    tb = last.shapes.add_textbox(Emu(457200), Emu(6400800), Emu(11278200), Emu(400000))
    p = tb.text_frame.paragraphs[0]
    p.text = "Note: PPTX designed for Cormorant Garamond. Install from Google Fonts for full fidelity."
    for run in p.runs:
        run.font.size = Pt(10)
        run.font.name = "IBM Plex Sans"
    prs.save(pptx_path)
    return True


def inline_generate(content_path: Path, design_spec_path: Path, template_path: Path | None, output_path: Path) -> None:
    content = yaml.safe_load(content_path.read_text()) or {}
    spec = yaml.safe_load(design_spec_path.read_text()) or {} if design_spec_path.exists() else {}

    # Always start from a fresh Presentation(). Loading template.pptx here was
    # producing duplicate `ppt/slides/slideN.xml` zip entries in the output —
    # symptom of stale slide-part paths inside the template colliding with new
    # add_slide() calls. A blank Presentation is deterministic and we don't
    # lose anything meaningful (the template had no slides to begin with; its
    # purpose was theme inheritance, which we re-apply explicitly via the
    # design spec's colors/fonts).
    prs = Presentation()
    prs.slide_width = Emu(spec.get("slide_dimensions", {}).get("width_emu", 12192000))
    prs.slide_height = Emu(spec.get("slide_dimensions", {}).get("height_emu", 6858000))

    theme = spec.get("theme_colors") or {}
    coral_hex = (theme.get("coral") or {}).get("rgb_hex", "FF795E")
    ink_hex   = (theme.get("ink")   or {}).get("rgb_hex", "161616")
    white_hex = (theme.get("white") or {}).get("rgb_hex", "FFFFFF")
    warm_hex  = (theme.get("warm")  or {}).get("rgb_hex", "F7F3EE")

    def solid(shape, hex_: str) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(hex_.upper())

    bg_for_layout = {
        "cover_layout": ink_hex,
        "chapter_divider_layout": ink_hex,
        "photo_break_layout": ink_hex,
        "metric_grid_layout": warm_hex,
        "portfolio_grid_layout": "F4F4F4",
    }
    blank = prs.slide_layouts[6]

    # Layouts that want a single full-bleed image behind text. For portfolio-grid
    # / product-detail / metric-grid the inline fallback skips images — those
    # layouts need card-level composition that this simplified generator doesn't
    # model. The PDF path is authoritative for those pages; PPTX users can edit
    # them manually or invoke the richer data-to-pptx.py generator if the schema
    # is ever reconciled.
    FULL_BLEED_LAYOUTS = {"cover_layout", "chapter_divider_layout", "photo_break_layout"}

    for slide_spec in content.get("slides") or []:
        slide = prs.slides.add_slide(blank)
        layout_name = slide_spec.get("layout_name")
        # Background
        bg = bg_for_layout.get(layout_name, white_hex)
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        solid(bg_shape, bg)
        bg_shape.line.fill.background()

        # Full-bleed image BEFORE text shapes, so text renders on top in z-order.
        # python-pptx preserves insertion order as render order, so adding the
        # image here (right after the background rectangle) puts it behind every
        # subsequent textbox. Cover/divider/photo-break templates expect this.
        images = slide_spec.get("images") or {}
        if layout_name in FULL_BLEED_LAYOUTS and images:
            first_img = next(iter(images.values()), None)
            if first_img:
                try:
                    ip = Path(first_img)
                    if ip.exists():
                        slide.shapes.add_picture(str(ip), 0, 0, prs.slide_width, prs.slide_height)
                except Exception as e:
                    print(f"  [warn] full-bleed image insert failed on slide {slide_spec.get('slide_number')}: {e}", file=sys.stderr)

        shape_content = slide_spec.get("shape_content") or {}
        is_dark = bg in (ink_hex, "104473", "465D53") or layout_name in FULL_BLEED_LAYOUTS
        text_color_hex = white_hex if is_dark else ink_hex

        # Eyebrow (coral mono)
        eyebrow = shape_content.get("eyebrow") or ""
        if eyebrow:
            eb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(11), Inches(0.3))
            eb.text_frame.text = eyebrow
            for run in eb.text_frame.paragraphs[0].runs:
                run.font.name = "IBM Plex Mono"; run.font.size = Pt(11); run.font.bold = True
                run.font.color.rgb = RGBColor.from_string(coral_hex.upper())

        # Title (Cormorant)
        title_lead = shape_content.get("title_lead") or ""
        accent = shape_content.get("italic_accent_phrase") or ""
        if title_lead or accent:
            top_y = Inches(2.2) if slide_spec.get("layout_name") in ("cover_layout", "chapter_divider_layout") else Inches(1.1)
            tb = slide.shapes.add_textbox(Inches(0.6), top_y, Inches(11), Inches(2.0))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            if title_lead:
                r1 = p.add_run(); r1.text = title_lead + (" " if accent else "")
                r1.font.name = "Cormorant Garamond"; r1.font.size = Pt(56); r1.font.color.rgb = RGBColor.from_string(text_color_hex.upper())
            if accent:
                r2 = p.add_run(); r2.text = accent
                r2.font.name = "Cormorant Garamond"; r2.font.italic = True; r2.font.size = Pt(56); r2.font.color.rgb = RGBColor.from_string(text_color_hex.upper())

        # Body (first body-like field)
        body_text = (shape_content.get("body_col_1") or shape_content.get("subtitle")
                     or shape_content.get("panel_body") or shape_content.get("body") or "")
        if body_text:
            bb = slide.shapes.add_textbox(Inches(0.6), Inches(3.3), Inches(11), Inches(3.0))
            bb.text_frame.text = body_text
            bb.text_frame.word_wrap = True
            for para in bb.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "IBM Plex Sans"; run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor.from_string(text_color_hex.upper())

        # Page number footer
        pn = shape_content.get("page_number_roman")
        total = shape_content.get("total_pages_roman")
        if pn and total:
            pf = slide.shapes.add_textbox(Inches(11), Inches(7.0), Inches(2), Inches(0.3))
            pf.text_frame.text = f"{pn} / {total}"
            for run in pf.text_frame.paragraphs[0].runs:
                run.font.name = "IBM Plex Mono"; run.font.size = Pt(10)
                run.font.color.rgb = RGBColor.from_string(("C6C6C6" if is_dark else "525252"))

        # (Full-bleed image already inserted above for cover/divider/photo-break.
        # Card-based layouts like portfolio-grid/product-detail are intentionally
        # not populated here — see FULL_BLEED_LAYOUTS comment.)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--pptx")
    ap.add_argument("--design-spec")
    ap.add_argument("--inline-generate", action="store_true")
    ap.add_argument("--content")
    ap.add_argument("--template")
    ap.add_argument("--output")
    args = ap.parse_args()

    if args.inline_generate:
        if not (args.content and args.output):
            sys.stderr.write("ERROR: --inline-generate needs --content and --output\n")
            return 1
        inline_generate(Path(args.content), Path(args.design_spec) if args.design_spec else Path(""), Path(args.template) if args.template else None, Path(args.output))
        return 0

    if not args.pptx:
        sys.stderr.write("ERROR: --pptx required\n")
        return 1
    pptx = Path(args.pptx)
    spec = yaml.safe_load(Path(args.design_spec).read_text()) if args.design_spec and Path(args.design_spec).exists() else None

    stripped = strip_theme_shadows(pptx)
    bg_count = enforce_explicit_backgrounds(pptx, spec)
    try_embed_fonts(pptx, spec)

    print(f"Cleanup: shadows stripped={stripped}, background fills applied to {bg_count} text shapes, font advisory added", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
