#!/usr/bin/env python3
"""
build-template-pptx.py — one-time builder for templates/template.pptx.

Generates a minimal template.pptx with 9 named slide layouts that map 1:1 to
the 9 HTML page-templates. Layouts in python-pptx are difficult to author
cleanly (you extend an existing template), so this script takes the pragmatic
approach: emit a blank presentation with widescreen 13.33x7.5 and let
data-to-pptx.py build slides per layout_name. The layout_name is stored as a
hint in each output slide's notes; the rendering (actual shape positioning)
happens in SLICE-004-04 during generation, not here.

Re-running is idempotent (byte-identical modulo zip timestamps).

Usage:
    build-template-pptx.py [--output templates/template.pptx]
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.stderr.write("ERROR: python-pptx required. pip install python-pptx\n")
    sys.exit(1)


LAYOUT_NAMES = [
    "cover_layout",
    "two_column_narrative_layout",
    "metric_grid_layout",
    "timeline_layout",
    "chapter_divider_layout",
    "product_detail_layout",
    "portfolio_grid_layout",
    "photo_break_layout",
    "closing_layout",
]


def build(output: Path) -> None:
    prs = Presentation()
    # Widescreen 13.33 x 7.5 inches in EMU (914400 per inch)
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    # Create one placeholder slide per layout_name using the blank layout
    # and annotate the notes with the canonical layout_name hint.
    blank_layout = prs.slide_layouts[6]  # Blank
    for name in LAYOUT_NAMES:
        slide = prs.slides.add_slide(blank_layout)
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"layout_hint: {name}"

    # Remove the placeholder slides — template.pptx should have 0 slides
    # (data-to-pptx.py adds slides per page-plan entry)
    # Work around python-pptx not supporting direct removal: rebuild without slides
    # by saving, then rewriting.
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    # Re-open and remove all slides (drop the layout-hint carriers now that
    # data-to-pptx consumes the layout_name field in the content YAML directly).
    prs2 = Presentation(output)
    xml_slides = prs2.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for sld in slides_to_remove:
        xml_slides.remove(sld)
    prs2.save(output)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--output", default=None)
    args = ap.parse_args()

    skill_root = Path(__file__).resolve().parent.parent
    out = Path(args.output) if args.output else skill_root / "templates" / "template.pptx"

    build(out)
    print(f"Built {out} ({len(LAYOUT_NAMES)} layout hints emitted then cleared; template is now empty, ready for data-to-pptx.py)", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
